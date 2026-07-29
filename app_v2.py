"""
Lobby-Ops Weekly Report Dashboard V2
- Auto-reads Major Updates.xlsx from workspace
- CSV uploads for W1/W2 game data
- Improved slide layout with centered region cards
- Larger graphs
- Weekly Summary tab: one-page summary with PDF/image export
"""
import streamlit as st
import csv, re, io, os, tempfile
from datetime import datetime
from collections import defaultdict, OrderedDict
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# === COLORS ===
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_LIGHT = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_HEADER = RGBColor(0x4C, 0x1D, 0x95)
DARK_TEXT = RGBColor(0x1F, 0x2A, 0x37)
GRAY_TEXT = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
PINK = RGBColor(0xD9, 0x66, 0xFF)
GOLD = RGBColor(0xFF, 0xAA, 0x01)
SILVER = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x05, 0x96, 0x69)
GOTM_CLR = RGBColor(0xBD, 0xF2, 0x41)
FONT = "Mulish"
SILVER_URL = "https://coralracing-my.sharepoint.com/:x:/g/personal/ssravya_entainindia_com/IQAc8F0vkrx7TJcoTNCHn99EAUy7zqknO5bhDP822a9_YUQ?e=R3LpzE"
WORKSPACE = os.path.dirname(os.path.abspath(__file__))
EXCEL_PATH = os.path.join(WORKSPACE, "Major Updates - Detailed Summary Report.xlsx")

MARKET_TO_SLIDE = {
    "UK": "UK", "Brazil": "Brazil", "Ontario - Entain": "Ontario",
    ".com (CA)": "ROC", ".com (EU)": "ROW", "Greece": "Greece",
    "Spain": "Southern Europe", "Portugal": "Southern Europe",
    "Colombia": "Southern Europe", "Denmark": "Denmark",
    "Germany": "Germany", "Austria": "Austria",
    "Belgium": "Belgium", "Italy": "Italy", "South Africa": "South Africa",
}

# Map Excel brand names to our slide keys
BRAND_TO_SLIDE = {
    "UK": "UK", "UK(X)": "UK", "UK-X(GAMING)": "UK",
    "UK(GAMING)": "UK", "UK GAMING": "UK",
    "BRAZIL": "Brazil", "Brazil": "Brazil",
    "CANADA(ROC)": "ROC", "CANADA & ONTARIO": "Ontario",
    "CANADA(Ontario)": "Ontario", "CANADA": "ROC",
    "Baltics": "ROW", "ROW": "ROW",
    "GREECE": "Greece", "Greece": "Greece",
    "SOUTHERN EUROPE": "Southern Europe",
    "Spain": "Southern Europe", "Portugal": "Southern Europe",
    "Colombia": "Southern Europe",
    "DENMARK": "Denmark", "Denmark": "Denmark",
    "GERMANY": "Germany", "Germany": "Germany",
    "AUSTRIA": "Austria", "Austria": "Austria",
    "BELGIUM": "Belgium", "Belgium": "Belgium",
    "ITALY": "Italy", "Italy": "Italy",
    "SOUTH AFRICA": "South Africa", "South Africa": "South Africa",
    "BIG": "Belgium",
}

REGION_ORDER = [
    ("UK Region", ["UK"], "Ladbrokes, Coral, Bwin.UK, Sportingbet.UK, Gamebookers.UK, Partycasino.UK, Partypoker.UK, Gala, Gala Casino, Foxy"),
    ("ASE Region", ["Brazil", "Ontario", "ROC", "Greece", "Southern Europe"], "Brazil, Ontario, ROC, Greece, Spain, Portugal, Colombia"),
    ("NCE Region", ["ROW", "Austria", "Denmark", "Germany"], "ROW, Austria, Denmark, Germany"),
    ("Others", ["Belgium", "Italy", "South Africa"], "Belgium, Italy, South Africa"),
]

# === EXCEL PARSER ===
def parse_excel(path):
    """Parse Major Updates.xlsx - returns game counts, optimizations, containers."""
    wb = openpyxl.load_workbook(path, data_only=True)

    # Sheet 1: Game counts
    ws1 = wb[wb.sheetnames[0]]
    game_counts = {}
    opt_counts = {}
    for r in range(2, ws1.max_row + 1):
        market = ws1.cell(r, 1).value
        games = ws1.cell(r, 4).value or ws1.cell(r, 2).value
        if market and games:
            game_counts[market.strip()] = int(games) if games else 0

    # Read optimization counts from Sheet 1 cols G-P
    # H=UK brands, J=S.Europe, L=America, N=NCE, P=Others
    opt_nums = {}  # slide_key -> count

    # UK column (G-H) - split into UK(X), Party, Gala but also provide full UK total
    uk_all = 0; party_opt = 0; gala_opt = 0; uk_x_opt = 0
    for r in range(2, ws1.max_row + 1):
        brand = ws1.cell(r, 7).value
        cnt = ws1.cell(r, 8).value
        if brand and cnt:
            bl = brand.strip().lower()
            try:
                cv = int(cnt)
            except (ValueError, TypeError):
                continue
            uk_all += cv
            if any(k in bl for k in ["party casino", "party poker"]):
                party_opt += cv
            elif any(k in bl for k in ["foxy", "gala"]):
                gala_opt += cv
            else:
                uk_x_opt += cv
    # UK region slide gets ALL brands combined
    opt_nums["UK"] = uk_all
    # Individual UK slides get their split
    if party_opt > 0: opt_nums["UK Gaming (Party)"] = party_opt
    if gala_opt > 0: opt_nums["UK Gaming (Gala/Foxy)"] = gala_opt

    # Southern Europe column (I-J): Brazil=sportingbet br+betboo, Greece=gr+vistabet, S.Europe=rest
    brazil_opt = 0; greece_opt = 0; seur_opt = 0
    for r in range(2, ws1.max_row + 1):
        brand = ws1.cell(r, 9).value
        cnt = ws1.cell(r, 10).value
        if brand and cnt:
            bl = brand.strip().lower()
            try:
                cv = int(cnt)
            except (ValueError, TypeError):
                continue
            if "sportingbet br" in bl or "betboo" in bl: brazil_opt += cv
            elif "gr" in bl or "vistabet" in bl: greece_opt += cv
            else: seur_opt += cv
    if brazil_opt > 0: opt_nums["Brazil"] = brazil_opt
    if greece_opt > 0: opt_nums["Greece"] = greece_opt
    if seur_opt > 0: opt_nums["Southern Europe"] = seur_opt

    # America column (K-L): ROC = non-Ontario brands, Ontario = Ontario brands
    roc_opt = 0; ont_opt = 0
    for r in range(2, ws1.max_row + 1):
        brand = ws1.cell(r, 11).value
        cnt = ws1.cell(r, 12).value
        if brand and cnt:
            bl = brand.strip().lower()
            try:
                cv = int(cnt)
            except (ValueError, TypeError):
                continue
            if "ontario" in bl: ont_opt += cv
            else: roc_opt += cv
    if roc_opt > 0: opt_nums["ROC"] = roc_opt
    if ont_opt > 0: opt_nums["Ontario"] = ont_opt

    # NCE column (M-N): ROW, Austria, Denmark, Germany
    row_opt = 0; austria_opt = 0; denmark_opt = 0; germany_opt = 0
    for r in range(2, ws1.max_row + 1):
        brand = ws1.cell(r, 13).value
        cnt = ws1.cell(r, 14).value
        if brand and cnt:
            bl = brand.strip().lower()
            try:
                cv = int(cnt)
            except (ValueError, TypeError):
                continue
            if "dk" in bl: denmark_opt += cv
            elif "(at)" in bl: austria_opt += cv
            elif " de" in bl or "bpremium de" in bl: germany_opt += cv
            else: row_opt += cv  # Bwin.com, Premium, Party Casino com, etc.
    if row_opt > 0: opt_nums["ROW"] = row_opt
    if austria_opt > 0: opt_nums["Austria"] = austria_opt
    if denmark_opt > 0: opt_nums["Denmark"] = denmark_opt
    if germany_opt > 0: opt_nums["Germany"] = germany_opt

    # Others column (O-P): Belgium, Italy, South Africa
    belgium_opt = 0; italy_opt = 0; sa_opt = 0
    for r in range(2, ws1.max_row + 1):
        brand = ws1.cell(r, 15).value
        cnt = ws1.cell(r, 16).value
        if brand and cnt:
            bl = brand.strip().lower()
            try:
                cv = int(cnt)
            except (ValueError, TypeError):
                continue
            if "be" in bl: belgium_opt += cv
            elif "it" in bl or "giocodigitale" in bl: italy_opt += cv
            elif "za" in bl: sa_opt += cv
    if belgium_opt > 0: opt_nums["Belgium"] = belgium_opt
    if italy_opt > 0: opt_nums["Italy"] = italy_opt
    if sa_opt > 0: opt_nums["South Africa"] = sa_opt

    # Sheet 2: Optimisations & Containers
    ws2 = wb[wb.sheetnames[1]]
    optimizations = defaultdict(list)
    containers = defaultdict(list)

    for r in range(2, ws2.max_row + 1):
        brand = ws2.cell(r, 1).value
        if not brand:
            continue
        brand = brand.strip()
        new_cont = ws2.cell(r, 3).value or ""
        other_upd = ws2.cell(r, 4).value or ""

        # Map brand to slide key
        slide_key = BRAND_TO_SLIDE.get(brand, brand)

        # Parse containers (bullet-separated)
        if new_cont and new_cont.strip() != "NA":
            for line in new_cont.split("\n"):
                line = line.strip().lstrip("•").lstrip("-").strip()
                if line and line != "NA":
                    containers[slide_key].append(line)

        # Parse optimizations (bullet-separated)
        if other_upd and other_upd.strip() != "NA":
            for line in other_upd.split("\n"):
                line = line.strip().lstrip("•").lstrip("-").strip()
                if line and line != "NA":
                    optimizations[slide_key].append(line)

    wb.close()
    return game_counts, optimizations, containers, opt_nums

# === CSV PARSER ===
def parse_date(date_str):
    date_str = date_str.strip()
    if not date_str:
        return None
    date_str = re.sub(r'(\d+)(st|nd|rd|th)', r'\1', date_str)
    parts = date_str.replace(',', '').split()
    months = ["January","February","March","April","May","June",
              "July","August","September","October","November","December"]
    month = day = year = None
    for i, p in enumerate(parts):
        if p in months:
            month = p
            if i+1 < len(parts) and parts[i+1].isdigit():
                day = parts[i+1]
            if i+2 < len(parts) and parts[i+2].isdigit() and len(parts[i+2]) == 4:
                year = parts[i+2]
            break
    if month and day and year:
        try:
            return datetime.strptime(f"{month} {day} {year}", "%B %d %Y")
        except:
            return None
    return None

def parse_csv(content):
    games = []
    reader = csv.DictReader(io.StringIO(content))
    headers = reader.fieldnames
    rating_col = next((h for h in headers if "Rating" in h), None)
    date_col = next((h for h in headers if "Published Date" in h), None)
    for row in reader:
        market = row.get("Market/Country (labels)", "").strip("[]")
        name = row.get("Task Name", "").strip()
        rating = row.get(rating_col, "").strip() if rating_col else ""
        pub_str = row.get(date_col, "").strip() if date_col else ""
        provider = row.get("Provider (labels)", "").strip("[]")
        pd = parse_date(pub_str)
        if name and market:
            # Treat "not seen" or empty rating as Silver
            if not rating or rating.lower() == "not seen":
                rating = "Silver"
            games.append({"name": name, "market": market, "rating": rating,
                          "pub_date": pd, "provider": provider})
    return games

def to_camel(name):
    words = name.strip().split()
    result = []
    for w in words:
        if w.isupper() and len(w) <= 4:
            result.append(w)
        elif w and w[0].isdigit():
            result.append(w)
        else:
            result.append(w.capitalize())
    return " ".join(result)

def group_by_market(games):
    grouped = defaultdict(lambda: defaultdict(list))
    for g in games:
        sk = MARKET_TO_SLIDE.get(g["market"])
        if not sk:
            continue
        display = f"{to_camel(g['name'])} by {g['provider']}" if g['provider'] else to_camel(g['name'])
        if display not in grouped[sk][g["rating"]]:
            grouped[sk][g["rating"]].append(display)
    return grouped

def daily_counts(games, market_key):
    counts = defaultdict(lambda: defaultdict(int))
    for g in games:
        if MARKET_TO_SLIDE.get(g["market"]) == market_key and g["pub_date"]:
            day = g["pub_date"].strftime("%a")
            counts[day][g["rating"]] += 1
    return counts

# === PPT HELPERS ===
def set_bg(slide):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = WHITE

def txt(slide, l, t, w, h, text, sz=11, bold=False, clr=DARK_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = clr; p.font.name = FONT; p.alignment = align
    return tb

def card(slide, l, t, w, h, label, value, accent=PURPLE, dark_bg=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if dark_bg:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x2D, 0x1B, 0x4E)
        sh.line.color.rgb = accent; sh.line.width = Pt(2)
        label_clr = RGBColor(0xC4, 0xB5, 0xFD)
        value_clr = WHITE
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = accent; sh.line.width = Pt(2)
        label_clr = GRAY_TEXT
        value_clr = DARK_TEXT
    tf = sh.text_frame; tf.margin_top = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(9); p.font.color.rgb = label_clr
    p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = str(value)
    p2.font.size = Pt(28); p2.font.bold = True
    p2.font.color.rgb = value_clr; p2.font.name = FONT; p2.alignment = PP_ALIGN.CENTER

def bullets_box(slide, l, t, w, h, items):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {item}"; p.font.size = Pt(10)
        p.font.color.rgb = GRAY_TEXT; p.font.name = FONT; p.space_before = Pt(4)

def game_column(slide, l, t, w, title, games, title_clr, max_show=6, hyperlink=None):
    """Not used anymore - replaced by game_table"""
    pass

def game_table(slide, l, t, w, plat, gold, silv, gotm=None):
    """Display games in a table - auto-fit, purple border, no game count in header."""
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn

    # Build raw columns
    raw_cols = []
    if gotm:
        raw_cols.append(("GOTM", gotm, GOTM_CLR))
    if plat:
        raw_cols.append(("Platinum", plat, PINK))
    if gold:
        raw_cols.append(("Gold", gold, GOLD))
    if silv:
        raw_cols.append(("Silver", silv, SILVER))

    if not raw_cols:
        return t + Inches(0.5)

    # Split any column that's too long
    lengths = [len(c[1]) for c in raw_cols]
    median_len = sorted(lengths)[len(lengths) // 2]
    split_threshold = max(10, int(median_len * 1.5)) if median_len > 3 else 10

    cols = []
    for title, games, clr in raw_cols:
        if len(games) > split_threshold:
            mid = (len(games) + 1) // 2
            cols.append((title, games[:mid], clr, 1, len(games)))
            cols.append((title, games[mid:], clr, 2, len(games)))
        else:
            cols.append((title, games, clr, 0, len(games)))

    num_cols = len(cols)

    # Ensure table fits within slide width
    table_w = min(w, Inches(12.5))

    # Create table
    tbl_shape = slide.shapes.add_table(2, num_cols, l, t, table_w, Inches(3.2))
    tbl = tbl_shape.table

    # Column widths proportional
    total_items = max(sum(len(c[1]) for c in cols), 1)
    for i in range(num_cols):
        ratio = max(len(cols[i][1]), 2) / total_items
        tbl.columns[i].width = int(table_w * max(ratio, 0.1))

    # Set table border color to purple
    tbl_xml = tbl._tbl
    tblPr = tbl_xml.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = tbl_xml.makeelement(qn('a:tblPr'), {})
        tbl_xml.insert(0, tblPr)

    # Header row
    tbl.rows[0].height = Inches(0.3)
    i = 0
    while i < num_cols:
        title, games, clr, part, full_count = cols[i]
        cell = tbl.cell(0, i)
        cell.text = title
        cell.fill.solid()
        cell.fill.fore_color.rgb = clr
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = DARK_TEXT if clr in [GOLD, GOTM_CLR] else WHITE
        p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # If this is part 1 of a split, merge with part 2 header
        if part == 1 and i + 1 < num_cols and cols[i + 1][3] == 2:
            tbl.cell(0, i).merge(tbl.cell(0, i + 1))
            i += 2
        else:
            i += 1

    # Data row
    max_games = max(len(c[1]) for c in cols)
    tbl.rows[1].height = Inches(max(2.5, max_games * 0.15))

    for i, (title, games, clr, part, full_count) in enumerate(cols):
        cell = tbl.cell(1, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(4); tf.margin_left = Pt(4)

        start_num = 1
        if part == 2:
            start_num = full_count - len(games) + 1

        for j, game in enumerate(games):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"{start_num + j}. {game}"
            p.font.size = Pt(9); p.font.color.rgb = DARK_TEXT
            p.font.name = FONT; p.space_before = Pt(2)
        cell.vertical_anchor = MSO_ANCHOR.TOP

        # Add purple border to each cell
        tcPr = cell._tc.get_or_add_tcPr()
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            border = tcPr.makeelement(qn(border_name), {})
            border.set('w', '12700')  # 1pt
            solidFill = border.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '5B21B6'})
            solidFill.append(srgbClr)
            border.append(solidFill)
            tcPr.append(border)

    return t + Inches(3.5)

def add_chart(slide, l, t, w, h, w1_data, w2_data):
    chart_data = CategoryChartData()
    cats = ["Mon", "Tue", "Wed", "Thu", "Fri"]
    chart_data.categories = ["W1", "", "", "", "", "W2", "", "", "", ""]
    for rating in ["Platinum", "Gold", "Silver"]:
        vals = [w1_data.get(d, {}).get(rating, 0) for d in cats]
        vals += [w2_data.get(d, {}).get(rating, 0) for d in cats]
        chart_data.add_series(rating, vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, l, t, w, h, chart_data)
    chart = cf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = DARK_TEXT; chart.legend.font.name = FONT
    colors = [PINK, GOLD, SILVER]
    for i, clr in enumerate(colors):
        s = chart.plots[0].series[i]
        s.format.line.color.rgb = clr; s.format.line.width = Pt(2); s.smooth = True
    chart.value_axis.major_gridlines.format.line.fill.background()
    chart.value_axis.has_minor_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.font.color.rgb = GRAY_TEXT
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.number_format = '0'
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = GRAY_TEXT
    chart.category_axis.has_major_gridlines = False

# === PPTX GENERATOR ===
def generate_pptx(w1_games, w2_games, optimizations, containers, week_label, highlight_title="", highlight_desc="", opt_nums=None, game_counts=None):
    if not opt_nums: opt_nums = {}
    if not game_counts: game_counts = {}
    # Use reference PPTX as template
    template_path = os.path.join(WORKSPACE, "LobbyOps_Weekly_Report_June_Week_1ref.pptx")
    if not os.path.exists(template_path):
        template_path = os.path.join(WORKSPACE, "Entain theme.pptx")
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        # Remove all existing slides from template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        # Layout 0 = Title slide (basketball player image)
        # Layout 31 = Blank (subtle diagonal bg - for all other slides)
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[31]
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
        title_layout = prs.slide_layouts[6]
        content_layout = prs.slide_layouts[6]

    # SLIDE 1: Title - uses branded layout (already has Entain logo, image, tagline)
    s = prs.slides.add_slide(title_layout)
    # Clear placeholder text
    for ph in s.placeholders:
        ph.text = ""
    txt(s, Inches(0.8), Inches(2.5), Inches(7), Inches(0.8), "WEEKLY REPORT", sz=40, bold=False, clr=PURPLE_HEADER)
    txt(s, Inches(0.8), Inches(3.3), Inches(6), Inches(0.4), f"{week_label}  |  P&T Global Gaming Content", sz=14, clr=PURPLE)

    # SLIDE 2: This Week's Highlights (uses content layout - subtle diagonal bg)
    s = prs.slides.add_slide(content_layout)
    txt(s, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5), "THIS WEEK'S HIGHLIGHTS", sz=20, bold=True, clr=PURPLE_HEADER)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.75), Inches(3), Inches(0.03))
    ln.fill.solid(); ln.fill.fore_color.rgb = PURPLE; ln.line.fill.background()
    # Highlight card
    hl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0))
    hl.fill.solid(); hl.fill.fore_color.rgb = WHITE
    hl.line.color.rgb = PURPLE; hl.line.width = Pt(1.5)
    tf = hl.text_frame; tf.margin_top = Pt(14); tf.margin_left = Pt(18); tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = highlight_title if highlight_title else "This Week's Key Highlights"
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = PURPLE; p.font.name = FONT
    p2 = tf.add_paragraph()
    p2.text = highlight_desc if highlight_desc else ""
    p2.font.size = Pt(11); p2.font.color.rgb = GRAY_TEXT; p2.font.name = FONT; p2.space_before = Pt(6)

    w2_grouped = group_by_market(w2_games)
    page = 3

    # Build slides for each region
    for region_name, market_keys, subtitle in REGION_ORDER:
        # Use Excel col D game counts for region totals (source of truth)
        region_market_map = {
            "UK": ["UK (X)", "UK (Party)", "UK (Gaming)"],
            "Brazil": ["Brazil"],
            "Ontario": ["Ontario"],
            "ROC": ["Canada"],
            "Greece": ["Greece"],
            "Southern Europe": ["Spain", "Portugal", "Colombia"],
            "ROW": ["ROW"],
            "Austria": ["Austria"],
            "Denmark": ["Denmark"],
            "Germany": ["Germany"],
            "Belgium": ["Belgium"],
            "Italy": ["Italy"],
            "South Africa": ["South Africa"],
        }
        total_games = 0
        if excel_status:
            for mk in market_keys:
                excel_names = region_market_map.get(mk, [mk])
                for en in excel_names:
                    total_games += game_counts.get(en, 0)
        if total_games == 0:
            total_games = sum(len(v) for mk in market_keys for v in w2_grouped.get(mk, {}).values())
        total_opts = sum(opt_nums.get(mk, len(optimizations.get(mk, []))) for mk in market_keys)
        total_cont = sum(len(containers.get(mk, [])) for mk in market_keys)

        # REGION DIVIDER - uses content layout (subtle diagonal bg)
        s = prs.slides.add_slide(content_layout)
        for ph in s.placeholders:
            ph.text = ""
        txt(s, Inches(0), Inches(1.7), Inches(13.333), Inches(0.8),
            region_name.upper(), sz=38, bold=True, clr=PURPLE_HEADER, align=PP_ALIGN.CENTER)
        txt(s, Inches(0), Inches(2.5), Inches(13.333), Inches(0.4),
            subtitle, sz=11, clr=GRAY_TEXT, align=PP_ALIGN.CENTER)

        # KPI cards - centered row (skip cards with 0 value)
        kpi_items = []
        if total_games > 0: kpi_items.append(("GAMES RELEASED", str(total_games), RGBColor(0x29,0x00,0x7D)))
        if total_opts > 0: kpi_items.append(("OPTIMIZATIONS", str(total_opts), RGBColor(0x00,0xE6,0xC3)))
        if total_cont > 0: kpi_items.append(("NEW CONTAINERS", str(total_cont), RGBColor(0x88,0x46,0xE6)))

        card_w = Inches(2.5); card_h = Inches(1.2); gap = Inches(0.4)
        total_width = len(kpi_items) * card_w.inches + (len(kpi_items)-1) * gap.inches if kpi_items else 0
        start_x = (13.333 - total_width) / 2
        cards_y = Inches(3.3)
        for i, (lbl, val, accent) in enumerate(kpi_items):
            card(s, Inches(start_x + i*(card_w.inches + gap.inches)), cards_y, card_w, card_h, lbl, val, accent=accent)

        # New Containers for this region - below KPI cards
        all_region_containers = []
        for mk in market_keys:
            for c in containers.get(mk, []):
                all_region_containers.append(c)

        if all_region_containers:
            # Light background for containers
            cont_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.2), Inches(4.7), Inches(10.9), Inches(2.5))
            cont_bg.fill.solid()
            cont_bg.fill.fore_color.rgb = WHITE
            cont_bg.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
            cont_bg.line.width = Pt(1)

            txt(s, Inches(1.5), Inches(4.85), Inches(10), Inches(0.3),
                "NEW CONTAINERS ENABLED", sz=11, bold=True, clr=PURPLE_HEADER, align=PP_ALIGN.CENTER)
            ln_rc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.15), Inches(2.3), Inches(0.02))
            ln_rc.fill.solid(); ln_rc.fill.fore_color.rgb = PURPLE; ln_rc.line.fill.background()
            bullets_box(s, Inches(1.5), Inches(5.25), Inches(10), Inches(1.8), all_region_containers)

        page += 1

        # Individual market slides
        for mk in market_keys:
            mk_games = w2_grouped.get(mk, {})
            # Austria uses same game data as ROW
            if mk == "Austria" and not mk_games:
                mk_games = w2_grouped.get("ROW", {})
            opts = optimizations.get(mk, [])
            # Show slide if it has games OR optimizations
            if not mk_games and not opts and not opt_nums.get(mk, 0):
                continue
            plat = mk_games.get("Platinum", [])
            gold = mk_games.get("Gold", [])
            silv = mk_games.get("Silver", [])
            gotm = mk_games.get("GOTM", [])
            total = len(plat) + len(gold) + len(silv) + len(gotm)
            opts = optimizations.get(mk, [])

            s = prs.slides.add_slide(content_layout)
            # Top bar
            bar_m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
            bar_m.fill.solid(); bar_m.fill.fore_color.rgb = PURPLE; bar_m.line.fill.background()
            txt(s, Inches(0.4), Inches(0.15), Inches(12), Inches(0.35),
                f"{mk.upper()}", sz=13, bold=True, clr=PURPLE_HEADER)
            ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.02))
            ln.fill.solid(); ln.fill.fore_color.rgb = PURPLE; ln.line.fill.background()

            # KPI cards top
            opt_count = opt_nums.get(mk, len(opts)) if opt_nums else len(opts)
            card(s, Inches(0.4), Inches(0.7), Inches(1.9), Inches(0.9), "GAMES RELEASED", str(total), accent=RGBColor(0x29,0x00,0x7D))
            card(s, Inches(2.5), Inches(0.7), Inches(1.9), Inches(0.9), "OPTIMIZATIONS", str(opt_count), accent=RGBColor(0x00,0xE6,0xC3))

            # Chart - below KPI cards (left half)
            # Austria uses same data as ROW
            chart_mk = "ROW" if mk == "Austria" else mk
            w1d = daily_counts(w1_games, chart_mk)
            w2d = daily_counts(w2_games, chart_mk)
            add_chart(s, Inches(0.4), Inches(1.8), Inches(4.4), Inches(2.4), w1d, w2d)

            # Games table - right side (tabular form) - constrained to slide width
            table_bottom = game_table(s, Inches(5.0), Inches(0.7), Inches(8.0),
                                       plat, gold, silv, gotm if gotm else None)

            # Optimizations + Containers on SECOND slide if too much content
            # Calculate available space
            content_start = max(Inches(4.4), table_bottom + Inches(0.2))

            # Optimizations - check if it fits, else new slide
            content_start = max(Inches(4.4), table_bottom + Inches(0.2))

            # If content would overflow slide (7.5 inches), create new slide for optimizations
            if content_start > Inches(5.5) and opts:
                # New slide for optimizations
                txt(s, Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.25), str(page), sz=8, clr=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
                page += 1
                s = prs.slides.add_slide(content_layout)
                bar_o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.04))
                bar_o.fill.solid(); bar_o.fill.fore_color.rgb = PURPLE; bar_o.line.fill.background()
                txt(s, Inches(0.4), Inches(0.15), Inches(12), Inches(0.35),
                    f"{mk.upper()} - CONTINUED", sz=13, bold=True, clr=PURPLE_HEADER)
                content_start = Inches(0.7)
                opt_width = Inches(12.5)
            elif content_start > Inches(4.6):
                # Fits but tight - constrain to graph width
                opt_width = Inches(4.8)
            else:
                opt_width = Inches(12.5)

            txt(s, Inches(0.4), content_start, opt_width, Inches(0.3), "OPTIMIZATIONS", sz=11, bold=True, clr=PURPLE_HEADER)
            ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), content_start + Inches(0.3), Inches(2), Inches(0.02))
            ln2.fill.solid(); ln2.fill.fore_color.rgb = PURPLE; ln2.line.fill.background()
            if opts:
                bullets_box(s, Inches(0.4), content_start + Inches(0.35), opt_width, Inches(2.5), opts)

            txt(s, Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.25), str(page), sz=8, clr=LIGHT_GRAY, align=PP_ALIGN.RIGHT)
            page += 1

    return prs

# === WEEKLY SUMMARY GENERATOR ===
SUMMARY_CARD_COLORS = [
    [RGBColor(0x6E, 0xE7, 0xB7), RGBColor(0xFD, 0xE6, 0x8A), RGBColor(0xFC, 0xA5, 0xA5),
     RGBColor(0xFD, 0xBA, 0x74), RGBColor(0x67, 0xE8, 0xF9)],
    [RGBColor(0xC4, 0xB5, 0xFD), RGBColor(0xF9, 0xA8, 0xD4), RGBColor(0xFC, 0xA5, 0xA5),
     RGBColor(0xFD, 0xBA, 0x74), RGBColor(0x5E, 0xEA, 0xD4)],
    [RGBColor(0xFD, 0xE6, 0x8A), RGBColor(0x6E, 0xE7, 0xB7), RGBColor(0x67, 0xE8, 0xF9),
     RGBColor(0xC4, 0xB5, 0xFD), RGBColor(0xF9, 0xA8, 0xD4)],
]

# Market display names for the summary grid
SUMMARY_MARKET_GRID_ORDER = [
    [("UK", "UK"), ("UK Gaming\n(Party)", "UK Gaming (Party)"), ("UK Gaming\n(Gala/Foxy)", "UK Gaming (Gala/Foxy)"), ("Brazil", "Brazil"), ("Canada", "Canada ROC")],
    [("Ontario", "Ontario"), ("Greece", "Greece"), ("S. Europe", "Southern Europe"), ("Austria", "Austria"), ("Denmark", "Denmark")],
    [("Germany", "Germany"), ("ROW", "ROW"), ("Belgium", "Belgium"), ("Italy", "Italy"), ("South Africa", "South Africa")],
]

# Map Excel Sheet1 market names to summary display names
EXCEL_MARKET_TO_SUMMARY = {
    # Sheet 1 exact names (game counts) - all UK variants map to single "UK"
    "UK": "UK",
    "UK (X)": "UK",
    "UK (Party)": "UK",
    "UK (Gaming)": "UK",
    "ROW": "ROW",
    "Austria": "Austria",
    "Denmark": "Denmark",
    "Spain": "Southern Europe",
    "Portugal": "Southern Europe",
    "Colombia": "Southern Europe",
    "Brazil": "Brazil",
    "Canada": "Canada ROC",
    "Ontario": "Ontario",
    "Belgium": "Belgium",
    "Italy": "Italy",
    "Greece": "Greece",
    "Germany": "Germany",
    "South Africa": "South Africa",
    "Bulgaria": "ROW",
    "Black Rush": "Black Rush",
    "Foxy NZ": "Foxy NZ",
    # Also handle variations that may appear
    "UK(X)": "UK",
    "UK(Gaming)": "UK",
    "UK-X(GAMING)": "UK",
    "UK(GAMING)": "UK",
    "UK GAMING": "UK",
    "BRAZIL": "Brazil",
    "CANADA(ROC)": "Canada ROC",
    "CANADA(Ontario)": "Ontario",
    "CANADA & ONTARIO": "Ontario",
    "GREECE": "Greece",
    "SOUTHERN EUROPE": "Southern Europe",
    "DENMARK": "Denmark",
    "GERMANY": "Germany",
    "AUSTRIA": "Austria",
    "BELGIUM": "Belgium",
    "ITALY": "Italy",
    "SOUTH AFRICA": "South Africa",
}

# Map for optimizations table display
EXCEL_BRAND_TO_OPT_DISPLAY = {
    "UK": "UK (Ladbrokes/Coral/Bwin/Sportingbet)",
    "UK-X(GAMING)": "UK Gaming (Partycasino/Partypoker)",
    "UK(GAMING)": "UK Gaming (Gala/Gala Casino/Foxy)",
    "UK(X)": "UK (Ladbrokes/Coral/Bwin/Sportingbet)",
    "BRAZIL": "Brazil",
    "Brazil": "Brazil",
    "CANADA(ROC)": "Canada ROC",
    "Canada(ROC)": "Canada ROC",
    "CANADA(Ontario)": "Ontario",
    "Canada(Ontario)": "Ontario",
    "CANADA & ONTARIO": "Ontario",
    "GREECE": "Greece",
    "Greece": "Greece",
    "SOUTHERN EUROPE": "Southern Europe",
    "Southern Europe": "Southern Europe",
    "DENMARK": "Denmark",
    "Denmark": "Denmark",
    "GERMANY": "Germany",
    "Germany": "Germany",
    "AUSTRIA": "Austria",
    "Austria": "Austria",
    "ROW": "ROW",
    "BELGIUM": "Belgium",
    "Belgium": "Belgium",
    "ITALY": "Italy",
    "Italy": "Italy",
    "SOUTH AFRICA": "South Africa",
    "South Africa": "South Africa",
}

def parse_pdf_for_summary(pdf_path):
    """Parse a summary PDF (like LobbyOps_Summary_Week5.pdf) to extract data.
    
    PDF text format (from pdfplumber):
    Line 0: LOBBY-OPS | WEEKLY SUMMARY
    Line 1: Week label
    Line 2: GAMES RELEASED OPTIMIZATIONS NEW CONTAINERS
    Line 3: 415 427 8
    Line 4: GAMES RELEASED BY MARKET
    Line 5: 36 28 34 44 18  (numbers row 1)
    Line 6: UK UK Gaming (Party) UK Gaming (Gala/Foxy) Brazil Canada ROC  (labels row 1)
    Line 7: 44 29 29 31 19  (numbers row 2)
    Line 8: Ontario Greece S. Europe Austria Denmark  (labels row 2)
    Line 9: 5 33 39 20 6  (numbers row 3)
    Line 10: Germany ROW Belgium Italy South Africa  (labels row 3)
    Then OPTIMIZATIONS BY MARKET / NEW CONTAINERS section
    Then KEY UPDATES section
    """
    import pdfplumber

    with pdfplumber.open(pdf_path) as pdf:
        page = pdf.pages[0]
        text = page.extract_text()

    lines = text.split('\n')

    # Parse game counts - look for lines with exactly 5 numbers
    game_counts = {}
    market_labels_order = [
        "UK", "UK Gaming (Party)", "UK Gaming (Gala/Foxy)", "Brazil", "Canada ROC",
        "Ontario", "Greece", "Southern Europe", "Austria", "Denmark",
        "Germany", "ROW", "Belgium", "Italy", "South Africa",
    ]

    numbers_rows = []
    for line in lines:
        parts = line.strip().split()
        if len(parts) == 5 and all(p.isdigit() for p in parts):
            numbers_rows.append([int(p) for p in parts])

    # First 3 rows of 5 numbers are the game counts grid
    idx = 0
    for row in numbers_rows[:3]:
        for count in row:
            if idx < len(market_labels_order):
                game_counts[market_labels_order[idx]] = count
                idx += 1

    # Parse optimizations - look for "MarketName Number" patterns
    opt_counts = {}
    known_opt_markets = [
        "UK (Ladbrokes/Coral/Bwin/Sportingbet)",
        "UK Gaming (Partycasino/Partypoker)",
        "UK Gaming (Gala/Gala Casino/Foxy)",
        "Brazil", "Canada ROC", "Ontario", "Greece",
        "Southern Europe", "Austria", "Denmark", "Germany",
        "ROW", "Belgium", "Italy", "South Africa",
    ]

    containers_raw = []
    in_opt_section = False

    for line in lines:
        if "OPTIMIZATIONS BY MARKET" in line or "Market Count" in line:
            in_opt_section = True
            continue
        if "KEY UPDATES" in line:
            break
        if not in_opt_section:
            continue

        line_s = line.strip()
        if not line_s or line_s == "Market Count" or line_s == "Container Brand":
            continue

        # Try to match optimization market + count
        matched_opt = False
        for market in known_opt_markets:
            if line_s.startswith(market):
                rest = line_s[len(market):].strip()
                num_match = re.match(r'^(\d+)', rest)
                if num_match:
                    opt_counts[market] = int(num_match.group(1))
                    # Anything after the number on same line is container info
                    container_rest = rest[num_match.end():].strip()
                    if container_rest:
                        _parse_container_line(container_rest, containers_raw)
                    matched_opt = True
                break

        if matched_opt:
            continue

        # Skip TOTAL line
        if line_s.startswith("TOTAL"):
            continue

        # Check if it's a container line (has brand-like text with dots)
        if re.search(r'\.\w{2}', line_s) or any(kw in line_s for kw in ["Gaming", "Casino", "Tournament", "Madness", "Slot", "Monopoly", "Eznav", "Football", "Pragmatic", "Merkur", "1x2"]):
            _parse_container_line(line_s, containers_raw)

    # Parse key updates
    key_updates = []
    in_key_section = False
    for line in lines:
        if "KEY UPDATES" in line:
            in_key_section = True
            continue
        if in_key_section:
            if "Entain" in line and "Lobby-Ops" in line:
                break
            line_s = line.strip()
            if line_s.startswith('\u2022'):
                key_updates.append(line_s.lstrip('\u2022').strip())
            elif line_s and key_updates:
                # Continuation of previous bullet
                key_updates[-1] += " " + line_s

    return game_counts, opt_counts, containers_raw, key_updates


def _parse_container_line(text, containers_list):
    """Parse container name and brand from a PDF text line."""
    text = text.strip()
    if not text or len(text) < 4:
        return
    # Skip if it's just a number
    if re.match(r'^\d+$', text):
        return

    # Common brand patterns (domain-based)
    brand_pattern = r'((?:Sportingbet|Betboo|Bwin|Partycasino|Partypoker|GD|Gala Casino|Gala|Ladbrokes|Coral)(?:\.\S+)?(?:\s*[,&]\s*(?:Sportingbet|Betboo|Bwin|Partycasino|Partypoker|GD|Gala Casino|Gala|Ladbrokes|Coral)(?:\.\S+)?)*)'

    m = re.search(brand_pattern, text)
    if m:
        brand = m.group(1).strip().rstrip(',')
        container_name = text[:m.start()].strip().rstrip(',').rstrip('&').strip()
        if container_name:
            containers_list.append((container_name, brand))
    else:
        # No brand found - store with empty brand
        containers_list.append((text, ""))


def parse_excel_for_summary(path):
    """Parse Major Updates.xlsx and extract all data needed for the summary page.
    
    Game counts: cols A & D (rows 2-19). Spain+Portugal+Colombia = Southern Europe.
    Optimisations: cols G-H (UK), I-J (Southern Europe/Brazil/Greece), K-L (America),
                   M-N (NCE), O-P (Others) - brand-level counts summed by region.
    """
    wb = openpyxl.load_workbook(path, data_only=True)
    ws1 = wb[wb.sheetnames[0]]

    # === GAME COUNTS from cols A & D (rows 2-19) ===
    game_counts = {}
    for r in range(2, 20):
        market = ws1.cell(r, 1).value
        total_games = ws1.cell(r, 4).value
        if not market or total_games is None:
            continue
        market = market.strip()
        try:
            count = int(total_games)
        except (ValueError, TypeError):
            continue
        # Rows 7,8,9 (Spain, Portugal, Colombia) sum into Southern Europe
        if market in ("Spain", "Portugal", "Colombia"):
            game_counts["Southern Europe"] = game_counts.get("Southern Europe", 0) + count
        else:
            # Map market name to summary display key
            key_map = {
                "UK": "UK", "ROW": "ROW", "Austria": "Austria", "Denmark": "Denmark",
                "Bulgaria": "ROW", "Brazil": "Brazil", "Canada": "Canada ROC",
                "Ontario": "Ontario", "Belgium": "Belgium", "Italy": "Italy",
                "Greece": "Greece", "Germany": "Germany", "South Africa": "South Africa",
                "Black Rush": "Black Rush", "Foxy NZ": "Foxy NZ",
            }
            summary_key = key_map.get(market)
            if summary_key:
                game_counts[summary_key] = game_counts.get(summary_key, 0) + count

    # === OPTIMISATIONS: read directly from G:H (Market | Count table) ===
    opt_counts = {}
    for r in range(2, ws1.max_row + 1):
        market = ws1.cell(r, 7).value
        cnt = ws1.cell(r, 8).value
        if market is None or str(market).strip() == "":
            continue
        market_clean = str(market).strip()
        if market_clean.upper() in ("TOTAL", "MARKET", "COUNT"):
            continue
        if cnt is None:
            opt_counts[market_clean] = 0
        else:
            try:
                opt_counts[market_clean] = int(cnt)
            except (ValueError, TypeError):
                opt_counts[market_clean] = 0

    # === CONTAINERS & KEY UPDATES from Sheet 2 ===
    ws2 = wb[wb.sheetnames[1]]
    all_containers = []
    key_updates = []

    for r in range(2, ws2.max_row + 1):
        brand = ws2.cell(r, 1).value
        if not brand:
            continue
        brand = brand.strip()
        new_cont = ws2.cell(r, 3).value or ""
        other_upd = ws2.cell(r, 4).value or ""

        # Map brand to display name for optimizations table
        opt_display = EXCEL_BRAND_TO_OPT_DISPLAY.get(brand, brand)

        # Parse containers
        if new_cont and new_cont.strip() != "NA":
            for line in new_cont.split("\n"):
                line = line.strip().lstrip("\u2022").lstrip("-").strip()
                if line and line != "NA":
                    all_containers.append((line, brand))

    # Build key updates from the most significant optimizations
    for r in range(2, ws2.max_row + 1):
        brand = ws2.cell(r, 1).value
        if not brand:
            continue
        brand = brand.strip()
        other_upd = ws2.cell(r, 4).value or ""
        if other_upd and other_upd.strip() != "NA":
            for line in other_upd.split("\n"):
                line = line.strip().lstrip("\u2022").lstrip("-").strip()
                if line and line != "NA" and len(line) > 30:
                    key_updates.append(line)

    wb.close()
    return game_counts, opt_counts, all_containers, key_updates


def generate_summary_pptx(game_counts, opt_counts, containers_raw, key_updates, week_label):
    """Generate the one-page weekly summary PPTX in PORTRAIT orientation."""
    import re as re_mod
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Portrait: 7.5 x 13.333 inches (swapped from landscape)
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.333)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    SW = prs.slide_width
    SH = prs.slide_height
    MARGIN = Emu(274320)
    CONTENT_W = SW - 2 * MARGIN

    # Colors
    PURPLE_ACCENT = RGBColor(0x5B, 0x21, 0xB6)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GRAY_S = RGBColor(0x1F, 0x1F, 0x1F)
    WHITE_S = RGBColor(0xFF, 0xFF, 0xFF)
    KPI_BG = RGBColor(0xF3, 0xF4, 0xF6)
    OPT_ROW_1 = RGBColor(0xED, 0xE9, 0xFE)
    OPT_ROW_2 = RGBColor(0xF5, 0xF3, 0xFF)
    TOTAL_ROW_BG = RGBColor(0xE8, 0xDE, 0xF8)

    def add_colored_rect(sl, left, top, width, height, fill_color):
        shape = sl.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_text(sl, left, top, width, height, text, font_size=Pt(11), bold=False,
                 color=DARK_GRAY_S, align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE):
        box = sl.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = v_anchor
        tf.margin_left = Emu(27432)
        tf.margin_right = Emu(27432)
        tf.margin_top = Emu(18288)
        tf.margin_bottom = Emu(18288)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = 'Mulish'
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def add_card_with_text(sl, left, top, width, height, fill_color, number, label):
        add_colored_rect(sl, left, top, width, height, fill_color)
        box = sl.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(18288)
        tf.margin_right = Emu(18288)
        tf.margin_top = Emu(18288)
        tf.margin_bottom = Emu(18288)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(number)
        run.font.name = 'Mulish'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY_S
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = label
        run2.font.name = 'Mulish'
        run2.font.size = Pt(7)
        run2.font.color.rgb = DARK_GRAY_S

    # Build games grid from parsed data
    games_grid = []
    for row_def in SUMMARY_MARKET_GRID_ORDER:
        row = []
        for card_label, lookup_key in row_def:
            count = game_counts.get(lookup_key, 0)
            row.append((card_label, count))
        games_grid.append(row)

    # Build optimizations list (ordered)
    opt_order = [
        "UK (Ladbrokes/Coral/Bwin/Sportingbet)",
        "UK Gaming (Partycasino/Partypoker)",
        "UK Gaming (Gala/Gala Casino/Foxy)",
        "Brazil", "Canada ROC", "Ontario", "Greece",
        "Southern Europe", "Austria", "Denmark", "Germany",
        "ROW", "Belgium", "Italy", "South Africa",
    ]
    optimizations_list = []
    for market in opt_order:
        if market in opt_counts:
            optimizations_list.append((market, opt_counts[market]))

    # Build containers list
    new_containers_list = []
    for desc, brand in containers_raw:
        match = re_mod.search(r"['\u2018\u2019]([^'\u2018\u2019]+)['\u2018\u2019]", desc)
        if match:
            container_name = match.group(1).strip()
        else:
            container_name = desc
        brand_display = EXCEL_BRAND_TO_OPT_DISPLAY.get(brand, brand)
        new_containers_list.append((container_name, brand_display))

    total_games = sum(c for row in games_grid for _, c in row)
    total_opts = sum(c for _, c in optimizations_list)
    total_containers = len(new_containers_list)

    # ============ PORTRAIT LAYOUT ============
    # Top accent bar
    add_colored_rect(slide, 0, 0, SW, Emu(54864), PURPLE_ACCENT)

    # Title
    add_text(slide, MARGIN, Emu(137160), CONTENT_W, Emu(365760),
             "LOBBY-OPS  |  WEEKLY SUMMARY", font_size=Pt(14), bold=True, color=PURPLE_ACCENT)
    add_text(slide, MARGIN, Emu(457200), CONTENT_W, Emu(228600),
             f"{week_label}  |  P&T Global Gaming Content", font_size=Pt(9), color=DARK_GRAY_S)

    # KPI Boxes - 3 across the width
    kpi_y = Emu(822960)
    kpi_h = Emu(594360)
    kpi_w = Emu(2011680)
    kpi_gap = Emu(137160)
    total_kpi_w = 3 * kpi_w + 2 * kpi_gap
    kpi_start_x = int((SW - total_kpi_w) / 2)

    kpis = [(total_games, "GAMES RELEASED"), (total_opts, "OPTIMIZATIONS"), (total_containers, "NEW CONTAINERS")]
    for idx, (val, lbl) in enumerate(kpis):
        x = kpi_start_x + idx * (kpi_w + kpi_gap)
        add_colored_rect(slide, x, kpi_y, kpi_w, kpi_h, KPI_BG)
        box = slide.shapes.add_textbox(x, kpi_y, kpi_w, kpi_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lbl
        run.font.name = 'Mulish'
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY_S
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = str(val)
        run2.font.name = 'Mulish'
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = PURPLE_ACCENT

    # GAMES RELEASED BY MARKET section
    section1_y = Emu(1600200)
    add_text(slide, MARGIN, section1_y, CONTENT_W, Emu(256032),
             "GAMES RELEASED BY MARKET", font_size=Pt(9), bold=True, color=PURPLE_ACCENT)
    add_colored_rect(slide, MARGIN, section1_y + Emu(256032), Emu(2286000), Emu(18288), PURPLE_ACCENT)

    # Colorful cards grid - 5 columns, 3 rows
    card_w = Emu(1188720)
    card_h = Emu(457200)
    card_gap_x = Emu(64008)
    card_gap_y = Emu(64008)
    total_cards_w = 5 * card_w + 4 * card_gap_x
    cards_start_x = int((SW - total_cards_w) / 2)
    grid_start_y = section1_y + Emu(320000)

    for row_idx, row in enumerate(games_grid):
        for col_idx, (label, count) in enumerate(row):
            x = cards_start_x + col_idx * (card_w + card_gap_x)
            y = grid_start_y + row_idx * (card_h + card_gap_y)
            color = SUMMARY_CARD_COLORS[row_idx][col_idx]
            add_card_with_text(slide, x, y, card_w, card_h, color, count, label)

    # OPTIMIZATIONS BY MARKET section
    opt_section_y = grid_start_y + 3 * (card_h + card_gap_y) + Emu(200000)
    add_text(slide, MARGIN, opt_section_y, Emu(3200000), Emu(256032),
             "OPTIMIZATIONS BY MARKET", font_size=Pt(9), bold=True, color=PURPLE_ACCENT)
    add_colored_rect(slide, MARGIN, opt_section_y + Emu(256032), Emu(2286000), Emu(18288), PURPLE_ACCENT)

    # Table header
    opt_header_y = opt_section_y + Emu(350000)
    row_h = Emu(155448)
    opt_col1_w = Emu(3200400)
    opt_col2_w = Emu(640080)

    add_colored_rect(slide, MARGIN, opt_header_y, opt_col1_w, Emu(9144), PURPLE_ACCENT)
    add_text(slide, MARGIN, opt_header_y, Emu(2514600), Emu(182880),
             "Market", font_size=Pt(9), bold=True, color=BLACK)
    add_text(slide, MARGIN + Emu(2560320), opt_header_y, opt_col2_w, Emu(182880),
             "Count", font_size=Pt(9), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Data rows
    opt_data = optimizations_list + [("TOTAL", total_opts)]
    data_start_y = opt_header_y + Emu(228600)

    for i, (market, count) in enumerate(opt_data):
        y = data_start_y + i * row_h
        is_total = (market == "TOTAL")
        bg = TOTAL_ROW_BG if is_total else (OPT_ROW_1 if i % 2 == 0 else OPT_ROW_2)
        add_colored_rect(slide, MARGIN, y, opt_col1_w, row_h, bg)
        add_text(slide, MARGIN + Emu(45720), y, Emu(2514600), row_h,
                 market, font_size=Pt(9), bold=is_total, color=BLACK)
        add_text(slide, MARGIN + Emu(2560320), y, opt_col2_w, row_h,
                 str(count), font_size=Pt(9), bold=is_total, color=BLACK, align=PP_ALIGN.CENTER)

    opt_table_end_y = data_start_y + len(opt_data) * row_h
    add_colored_rect(slide, MARGIN, opt_table_end_y, opt_col1_w, Emu(9144), PURPLE_ACCENT)

    # NEW CONTAINERS ENABLED section (below optimizations in portrait mode)
    cont_section_y = opt_table_end_y + Emu(250000)
    add_text(slide, MARGIN, cont_section_y, CONTENT_W, Emu(256032),
             "NEW CONTAINERS ENABLED", font_size=Pt(9), bold=True, color=PURPLE_ACCENT)
    add_colored_rect(slide, MARGIN, cont_section_y + Emu(256032), Emu(2286000), Emu(18288), PURPLE_ACCENT)

    # Container table - auto-fit columns
    cont_header_y = cont_section_y + Emu(350000)
    cont_row_h = Emu(180000)

    # Calculate column widths based on content length
    max_cont_name = max((len(c[0]) for c in new_containers_list), default=20)
    max_brand_name = max((len(c[1]) for c in new_containers_list), default=15)
    cont_total_w = int(CONTENT_W)
    cont_col1_w = int(cont_total_w * max_cont_name / (max_cont_name + max_brand_name))
    cont_col2_w = cont_total_w - cont_col1_w

    add_colored_rect(slide, MARGIN, cont_header_y, cont_total_w, Emu(9144), PURPLE_ACCENT)
    add_text(slide, MARGIN, cont_header_y, cont_col1_w, Emu(182880),
             "Container", font_size=Pt(9), bold=True, color=BLACK)
    add_text(slide, MARGIN + cont_col1_w, cont_header_y, cont_col2_w, Emu(182880),
             "Brand", font_size=Pt(9), bold=True, color=BLACK)

    cont_data_start_y = cont_header_y + Emu(228600)
    for i, (container, brand) in enumerate(new_containers_list):
        y = cont_data_start_y + i * cont_row_h
        bg = OPT_ROW_1 if i % 2 == 0 else OPT_ROW_2
        add_colored_rect(slide, MARGIN, y, cont_total_w, cont_row_h, bg)
        # Left-center aligned, no truncation
        add_text(slide, MARGIN + Emu(45720), y, cont_col1_w - Emu(45720), cont_row_h,
                 container, font_size=Pt(8), color=BLACK, align=PP_ALIGN.LEFT)
        add_text(slide, MARGIN + cont_col1_w, y, cont_col2_w, cont_row_h,
                 brand, font_size=Pt(8), color=BLACK, align=PP_ALIGN.LEFT)

    cont_table_end_y = cont_data_start_y + len(new_containers_list) * cont_row_h

    # KEY UPDATES section
    key_section_y = cont_table_end_y + Emu(250000)
    add_text(slide, MARGIN, key_section_y, CONTENT_W, Emu(256032),
             "KEY UPDATES", font_size=Pt(9), bold=True, color=PURPLE_ACCENT)
    add_colored_rect(slide, MARGIN, key_section_y + Emu(256032), Emu(2286000), Emu(18288), PURPLE_ACCENT)

    updates_y = key_section_y + Emu(320000)
    remaining = SH - updates_y - Emu(350000)

    # Auto-scale font size to fit all updates
    num_updates = len(key_updates)
    if num_updates > 12:
        update_font_size = Pt(5.5)
    elif num_updates > 8:
        update_font_size = Pt(6)
    else:
        update_font_size = Pt(6.5)

    box = slide.shapes.add_textbox(MARGIN, updates_y, CONTENT_W, remaining)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(36000)

    for idx, update in enumerate(key_updates):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = f"\u2022 {update}"
        run.font.name = 'Mulish'
        run.font.size = update_font_size
        run.font.color.rgb = BLACK

    # Footer
    add_text(slide, MARGIN, SH - Emu(280000), CONTENT_W, Emu(228600),
             "Entain  |  P&T Global Gaming Content  |  Lobby-Ops",
             font_size=Pt(7), color=PURPLE_ACCENT, align=PP_ALIGN.CENTER)
    add_colored_rect(slide, 0, SH - Emu(54864), SW, Emu(54864), PURPLE_ACCENT)

    return prs


# === STREAMLIT APP ===
st.set_page_config(page_title="Lobby-Ops Report Generator", page_icon="\U0001f3b0", layout="wide")

st.markdown("""
<style>
    .stApp {
        background-color: #f8f7fc;
    }
    [data-testid="stHeader"] { background-color: #f8f7fc; }
    .main-header { color: #4C1D95; font-size: 2.4rem; font-weight: 700; }
    .sub-header { color: #6B7280; font-size: 1rem; }
    h1, h2, h3 { color: #4C1D95 !important; }
    .stFileUploader label, .stTextInput label { color: #4C1D95 !important; }
    [data-testid="stFileUploader"] {
        border: 1px dashed #7c3aed; border-radius: 10px;
        padding: 10px; background: #faf8ff;
    }
    [data-testid="stFileUploader"] p, [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] small, [data-testid="stFileUploader"] div {
        color: #1f2a37 !important;
    }
    [data-testid="stFileUploadDropzone"] {
        background: #f5f3ff !important;
        color: #1f2a37 !important;
        border: none !important;
    }
    [data-testid="stFileUploadDropzone"] span {
        color: #4C1D95 !important;
        font-weight: 500 !important;
    }
    [data-testid="stFileUploadDropzone"] button {
        background: #7c3aed !important;
        color: white !important;
    }
    .stButton > button[kind="primary"] {
        background: linear-gradient(90deg, #7c3aed, #a855f7);
        color: white; border: none; border-radius: 8px;
        font-size: 1.1rem; font-weight: 600; padding: 0.6rem 2rem;
        width: 100%;
    }
    .stDownloadButton > button {
        background: linear-gradient(90deg, #7c3aed, #a855f7);
        color: white; border: none; border-radius: 8px;
        width: 100%; padding: 0.6rem 2rem; font-size: 1rem;
    }
    p, span, li, label { color: #1f2a37 !important; }
    .stTextInput input, .stNumberInput input, .stTextArea textarea {
        background-color: #ffffff; color: #1f2a37; border: 1px solid #e5e7eb;
    }
    .info-box { background: #f5f3ff; border: 1px solid #7c3aed; border-radius: 8px; padding: 12px; margin: 8px 0; color: #4C1D95; }
    .stAlert { width: 100%; }
    /* Force light mode on all dark elements */
    [data-testid="stAppViewContainer"], [data-testid="stSidebar"],
    [data-testid="stBottomBlockContainer"], .stTabs [data-baseweb="tab-panel"] {
        background-color: #f8f7fc !important;
    }
    [data-testid="stMetric"], [data-testid="stMetricValue"] {
        color: #1f2a37 !important;
    }
    [data-baseweb="tab"] {
        color: #4C1D95 !important;
        background-color: #f8f7fc !important;
    }
    [data-baseweb="input"], [data-baseweb="textarea"] {
        background-color: #ffffff !important;
        color: #1f2a37 !important;
    }
    div[data-testid="stNumberInput"] label,
    div[data-testid="stTextInput"] label,
    div[data-testid="stTextArea"] label,
    div[data-testid="stFileUploader"] label {
        color: #4C1D95 !important;
    }
    .stMarkdown, .stCaption { color: #1f2a37 !important; }
</style>
""", unsafe_allow_html=True)

st.markdown('<p class="main-header">\U0001f3b0 Lobby-Ops Weekly Report Generator</p>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">Entain P&T Global Gaming Content \u2014 Upload CSVs, generate the report in one click.</p>', unsafe_allow_html=True)

# Check if Excel exists
excel_status = os.path.exists(EXCEL_PATH)

# === TABS ===
tab1, tab2 = st.tabs(["\U0001f4ca Weekly Report", "\U0001f4dd Quick Summary"])

# ============ TAB 1: Weekly Report (existing) ============
with tab1:
    if excel_status:
        st.markdown(f'<div class="info-box">\u2705 <b>Major Updates.xlsx</b> detected in workspace (auto-loaded)</div>', unsafe_allow_html=True)
    else:
        st.warning("Major Updates.xlsx not found in workspace. Upload it below or place it in the Kiro workspace folder.")

    st.divider()

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("\U0001f4c4 W1 CSV (Previous Week)")
        w1_file = st.file_uploader("Upload W1 ClickUp CSV", type=["csv"], key="w1")
    with col2:
        st.subheader("\U0001f4c4 W2 CSV (Current Week)")
        w2_file = st.file_uploader("Upload W2 ClickUp CSV", type=["csv"], key="w2")

    st.divider()
    week_label = st.text_input("Week Label (e.g. 'May Week 4')", value="Jun Week 1", key="wl_report")

    st.subheader("Key Highlights of the Week")
    st.caption("This will appear on Slide 2 of the report. Enter title and description.")
    highlight_title = st.text_input("Highlight Title", value="", placeholder="e.g. Casino 2.0 Platform Launch - Gamebookers", key="hl_title")
    highlight_desc = st.text_area("Highlight Description", value="", placeholder="e.g. The new Casino 2.0 platform was successfully launched with 4,900+ games configured.", height=80, key="hl_desc")

    if st.button("\U0001f680 Generate Report", type="primary", use_container_width=True):
        if not w2_file:
            st.error("Please upload at least the W2 (current week) CSV.")
        else:
            with st.spinner("Generating your report..."):
                w1_games = []
                if w1_file:
                    w1_games = parse_csv(w1_file.read().decode("utf-8"))
                w2_games = parse_csv(w2_file.read().decode("utf-8"))
                optimizations = defaultdict(list)
                containers = defaultdict(list)
                if excel_status:
                    try:
                        game_counts, optimizations, containers, opt_nums = parse_excel(EXCEL_PATH)
                    except Exception as e:
                        st.warning(f"Excel parse issue: {e}. Proceeding without it.")
                        opt_nums = {}
                        game_counts = {}
                prs = generate_pptx(w1_games, w2_games, optimizations, containers, week_label, highlight_title, highlight_desc, opt_nums, game_counts)
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
                prs.save(tmp.name); tmp.close()
                st.success(f"\u2705 Report generated! W1: {len(w1_games)} games | W2: {len(w2_games)} games | Optimizations from Excel: {sum(len(v) for v in optimizations.values())} items")
                with open(tmp.name, "rb") as f:
                    pptx_bytes = f.read()
                col_a, col_b = st.columns(2)
                with col_a:
                    st.download_button(
                        label="\U0001f4e5 Download PPTX",
                        data=pptx_bytes,
                        file_name=f"LobbyOps_Weekly_Report_{week_label.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                with col_b:
                    st.download_button(
                        label="\U0001f4a1 PDF: Open PPTX \u2192 File \u2192 Export \u2192 PDF",
                        data=pptx_bytes,
                        file_name=f"LobbyOps_Weekly_Report_{week_label.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                os.unlink(tmp.name)

# ============ TAB 2: Weekly Summary ============
# Helper functions for Tab 2

def _parse_pptx_report(path):
    """Parse detailed weekly report PPTX - handles both old and new formats."""
    prs_in = Presentation(path)
    game_counts = {}
    opt_counts = {}
    containers = []
    key_updates = []

    # Market identification - match start of title
    market_patterns = [
        ("UK GAMING", "UK Gaming (Party)"),  # Will be refined below
        ("UK", "UK"),
        ("BRAZIL", "Brazil"),
        ("CANADA", "Canada ROC"),  # Will check for Ontario below
        ("ONTARIO", "Ontario"),
        ("GREECE", "Greece"),
        ("SOUTHERN EUROPE", "Southern Europe"),
        ("ROW", "ROW"),
        ("AUSTRIA", "Austria"),
        ("DENMARK", "Denmark"),
        ("GERMANY", "Germany"),
        ("BELGIUM", "Belgium"),
        ("ITALY", "Italy"),
        ("SOUTH AFRICA", "South Africa"),
    ]

    summary_to_opt = {
        "UK": "UK (Ladbrokes/Coral/Bwin/Sportingbet)",
        "UK Gaming (Party)": "UK Gaming (Partycasino/Partypoker)",
        "UK Gaming (Gala/Foxy)": "UK Gaming (Gala/Gala Casino/Foxy)",
        "Brazil": "Brazil", "Canada ROC": "Canada ROC", "Ontario": "Ontario",
        "Greece": "Greece", "Southern Europe": "Southern Europe",
        "ROW": "ROW", "Austria": "Austria", "Denmark": "Denmark",
        "Germany": "Germany", "Belgium": "Belgium", "Italy": "Italy",
        "South Africa": "South Africa",
    }

    for sl in prs_in.slides:
        texts = [s.text.strip() for s in sl.shapes if hasattr(s, "text") and s.text.strip()]
        if not texts:
            continue

        # Join all text for this slide to search
        all_text = "\n".join(texts)
        title = texts[0].upper().strip()

        # Skip title slide, region dividers, highlights
        if "LOBBY" in title and "OPS" in title:
            continue
        if title.endswith("REGION") or title.startswith("OTHERS"):
            # Region divider - look for containers
            for t in texts:
                if "container" in t.lower() and ("name" in t.lower() or "label" in t.lower()):
                    continue  # skip header row
                # Look for container table data
            continue
        if "HIGHLIGHT" in title or "WEEKLY REPORT" in title:
            continue

        # Identify market from title
        market_key = None
        for pattern, key in market_patterns:
            if title.startswith(pattern):
                # Special handling for UK GAMING
                if pattern == "UK GAMING":
                    if "PARTY" in title or "PARTYCASINO" in title:
                        market_key = "UK Gaming (Party)"
                    elif "GALA" in title or "FOXY" in title:
                        market_key = "UK Gaming (Gala/Foxy)"
                    else:
                        market_key = "UK Gaming (Party)"  # default
                elif pattern == "CANADA":
                    if "ONTARIO" in title or "(ONTARIO)" in title:
                        market_key = "Ontario"
                    else:
                        market_key = "Canada ROC"
                else:
                    market_key = key
                break

        if not market_key:
            continue

        # Extract Total Games Released and Total Optimizations
        games = 0
        opts = 0

        for i, t in enumerate(texts):
            t_upper = t.upper()
            # Look for "Total Games Released" followed by a number
            if "TOTAL GAMES RELEASED" in t_upper or "GAMES RELEASED" in t_upper:
                # Number might be in same text block or next text
                for p in t.split('\n'):
                    if p.strip().isdigit():
                        games = int(p.strip()); break
                # If not found in same block, check next text
                if games == 0 and i + 1 < len(texts) and texts[i+1].strip().isdigit():
                    games = int(texts[i+1].strip())

            if "TOTAL OPTIMIZATIONS" in t_upper or ("OPTIMIZATIONS" in t_upper and "BY MARKET" not in t_upper):
                for p in t.split('\n'):
                    if p.strip().isdigit():
                        opts = int(p.strip()); break
                if opts == 0 and i + 1 < len(texts) and texts[i+1].strip().isdigit():
                    opts = int(texts[i+1].strip())

            # Optimization bullets
            if "\u2022" in t or t.startswith("\u2022"):
                for line in t.split('\n'):
                    line = line.strip().lstrip("\u2022").lstrip("-").strip()
                    if line and len(line) > 20:
                        key_updates.append(line)

            # Container tables - look for "Container name" header then data
            if "CONTAINER" in t_upper and "NAME" in t_upper:
                # Next texts might be container data rows
                pass

        # Also look for standalone numbers that appear right after "Total Games Released"
        # by checking if any text is just a number and the previous text contains the label
        for i in range(1, len(texts)):
            if texts[i].strip().isdigit():
                prev = texts[i-1].upper()
                val = int(texts[i].strip())
                if "TOTAL GAMES RELEASED" in prev and games == 0:
                    games = val
                elif "TOTAL OPTIMIZATIONS" in prev and opts == 0:
                    opts = val

        if games > 0:
            game_counts[market_key] = games
        if opts > 0:
            opt_counts[summary_to_opt.get(market_key, market_key)] = opts

        # Extract containers from this slide
        # Look for patterns like "Container name | Label | Position" table
        in_container_section = False
        for t in texts:
            if "Container name" in t or "New Container" in t.replace("s", ""):
                in_container_section = True
                continue
            if in_container_section and t.strip() and not t.strip().isdigit():
                # Try to parse container row
                parts = t.split('\n')
                for part in parts:
                    part = part.strip()
                    if part and len(part) > 3 and not part.isdigit() and "Container" not in part and "Label" not in part and "Position" not in part:
                        # This might be a container entry
                        # Format could be: "ContainerName Label Position" on one line
                        # Or separate texts
                        containers.append((part, market_key))

    return game_counts, opt_counts, containers, key_updates


def _parse_summary_pdf(path):
    """Parse PDF - handles both summary (1-page) and detailed report (multi-page) formats."""
    import pdfplumber
    with pdfplumber.open(path) as pdf:
        all_text = ""
        page_texts = []
        for page in pdf.pages:
            pt = page.extract_text() or ""
            page_texts.append(pt)
            all_text += pt + "\n"

    lines = all_text.split('\n')

    # Detect format: if first page has "GAMES RELEASED BY MARKET" it's a summary
    is_summary = "GAMES RELEASED BY MARKET" in page_texts[0] if page_texts else False

    labels = ["UK", "UK Gaming (Party)", "UK Gaming (Gala/Foxy)", "Brazil", "Canada ROC",
              "Ontario", "Greece", "Southern Europe", "Austria", "Denmark",
              "Germany", "ROW", "Belgium", "Italy", "South Africa"]
    opt_order = [
        "UK (Ladbrokes/Coral/Bwin/Sportingbet)", "UK Gaming (Partycasino/Partypoker)",
        "UK Gaming (Gala/Gala Casino/Foxy)", "Brazil", "Canada ROC", "Ontario", "Greece",
        "Southern Europe", "Austria", "Denmark", "Germany", "ROW", "Belgium", "Italy", "South Africa",
    ]

    game_counts = {}
    opt_counts = {}
    containers = []
    key_updates = []

    if is_summary:
        # Parse one-page summary format
        num_rows = []
        for line in lines:
            parts = line.strip().split()
            if len(parts) == 5 and all(p.isdigit() for p in parts):
                num_rows.append([int(p) for p in parts])
        idx = 0
        for row in num_rows[:3]:
            for c in row:
                if idx < len(labels): game_counts[labels[idx]] = c; idx += 1

        in_opt = False
        for line in lines:
            if "OPTIMIZATIONS BY MARKET" in line or "Market Count" in line: in_opt = True; continue
            if "KEY UPDATES" in line: break
            if not in_opt: continue
            ls = line.strip()
            if not ls or ls in ("Market Count", "Container Brand"): continue
            for mkt in opt_order:
                if ls.startswith(mkt):
                    rest = ls[len(mkt):].strip()
                    m = re.match(r'^(\d+)', rest)
                    if m:
                        opt_counts[mkt] = int(m.group(1))
                        cr = rest[m.end():].strip()
                        if cr and len(cr) > 3:
                            bm = re.search(r'((?:Sportingbet|Betboo|Bwin|Partycasino|Gala Casino|GD)\S*)', cr)
                            if bm:
                                containers.append((cr[:bm.start()].strip().rstrip(','), cr[bm.start():].strip()))
                            else:
                                containers.append((cr, ""))
                    break

        in_ku = False
        for line in lines:
            if "KEY UPDATES" in line: in_ku = True; continue
            if in_ku:
                if "Entain" in line and "Lobby-Ops" in line: break
                ls = line.strip()
                if ls.startswith('\u2022'): key_updates.append(ls.lstrip('\u2022').strip())
                elif ls and key_updates: key_updates[-1] += " " + ls
    else:
        # Parse multi-page detailed report PDF
        # Look for patterns: "GAMES RELEASED\nN" and market names followed by data
        title_to_key = {
            "UK": "UK", "BRAZIL": "Brazil", "ONTARIO": "Ontario",
            "ROC": "Canada ROC", "GREECE": "Greece",
            "SOUTHERN EUROPE": "Southern Europe", "ROW": "ROW",
            "AUSTRIA": "Austria", "DENMARK": "Denmark", "GERMANY": "Germany",
            "BELGIUM": "Belgium", "ITALY": "Italy", "SOUTH AFRICA": "South Africa",
        }
        summary_to_opt = {
            "UK": "UK (Ladbrokes/Coral/Bwin/Sportingbet)", "Brazil": "Brazil",
            "Canada ROC": "Canada ROC", "Ontario": "Ontario", "Greece": "Greece",
            "Southern Europe": "Southern Europe", "ROW": "ROW", "Austria": "Austria",
            "Denmark": "Denmark", "Germany": "Germany", "Belgium": "Belgium",
            "Italy": "Italy", "South Africa": "South Africa",
        }

        # Process each page to find market data
        region_data = {}
        for page_text in page_texts:
            plines = page_text.split('\n')
            if not plines:
                continue

            # Check if this is a region divider or market detail page
            first_line = plines[0].strip().upper()

            # Region page
            if "REGION" in first_line or first_line == "OTHERS":
                games = opts = 0
                for pl in plines:
                    if "GAMES RELEASED" in pl:
                        for p in pl.split():
                            if p.isdigit(): games = int(p); break
                    if "OPTIMIZATIONS" in pl and "BY MARKET" not in pl:
                        for p in pl.split():
                            if p.isdigit(): opts = int(p); break
                    if "\u2022" in pl or pl.strip().startswith("•"):
                        bullet = pl.strip().lstrip("\u2022").lstrip("•").strip()
                        if bullet and len(bullet) > 5:
                            containers.append((bullet, first_line))
                region_data[first_line] = {"games": games, "opts": opts}
                continue

            # Market page
            clean = first_line.split("|")[0].strip().split("  ")[0].strip()
            if clean in title_to_key:
                sk = title_to_key[clean]
                games = opts = 0
                for pl in plines:
                    if "GAMES RELEASED" in pl:
                        for p in pl.split():
                            if p.isdigit(): games = int(p); break
                    if "OPTIMIZATIONS" in pl and "BY MARKET" not in pl:
                        for p in pl.split():
                            if p.isdigit(): opts = int(p); break
                    if "\u2022" in pl or pl.strip().startswith("•"):
                        bullet = pl.strip().lstrip("\u2022").lstrip("•").strip()
                        if bullet and len(bullet) > 20:
                            key_updates.append(bullet)
                if games > 0:
                    game_counts[sk] = games
                if opts > 0:
                    opt_counts[summary_to_opt.get(sk, sk)] = opts

        # Calculate Austria if missing
        if "Austria" not in game_counts and "NCE REGION" in region_data:
            nce = region_data["NCE REGION"]
            a_g = nce["games"] - game_counts.get("ROW", 0) - game_counts.get("Denmark", 0) - game_counts.get("Germany", 0)
            if a_g > 0:
                game_counts["Austria"] = a_g

    return game_counts, opt_counts, containers, key_updates




def _gen_summary_png(game_counts, opt_counts, cont_list, key_updates, week_label, total_games_override=None, total_opts_override=None, total_cont_override=None):
    """Generate summary report as PNG using Pillow - works on any platform."""
    from PIL import Image, ImageDraw, ImageFont
    import io as _io

    # Canvas settings
    W = 960  # width in pixels
    MARGIN = 40
    CW = W - 2 * MARGIN  # content width

    # Colors
    PURPLE_C = (91, 33, 182)
    DARK_C = (31, 31, 31)
    GRAY_C = (75, 85, 99)
    KPI_BG_C = (243, 244, 246)
    ROW1_C = (237, 233, 254)
    ROW2_C = (245, 243, 255)
    TOTAL_BG_C = (232, 222, 248)
    WHITE_C = (255, 255, 255)

    CARD_COLORS = [
        [(110,231,183),(253,230,138),(252,165,165),(253,186,116),(103,232,249)],
        [(196,181,253),(249,168,212),(252,165,165),(253,186,116),(94,234,212)],
        [(253,230,138),(110,231,183),(103,232,249),(196,181,253),(249,168,212)],
    ]

    # Try to load Mulish font, fallback to default
    try:
        font_dir = os.path.dirname(os.path.abspath(__file__))
        mulish_reg = os.path.join(font_dir, "Mulish-Regular.ttf")
        mulish_bold = os.path.join(font_dir, "Mulish-Bold.ttf")
        if not os.path.exists(mulish_reg):
            mulish_reg = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
            mulish_bold = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        font_sm = ImageFont.truetype(mulish_reg, 15)
        font_md = ImageFont.truetype(mulish_bold, 18)
        font_lg = ImageFont.truetype(mulish_bold, 28)
        font_xl = ImageFont.truetype(mulish_bold, 40)
        font_title = ImageFont.truetype(mulish_bold, 22)
        font_tiny = ImageFont.truetype(mulish_reg, 13)
        font_section = ImageFont.truetype(mulish_bold, 15)
    except:
        font_sm = ImageFont.load_default()
        font_md = font_sm
        font_lg = font_sm
        font_xl = font_sm
        font_title = font_sm
        font_tiny = font_sm
        font_section = font_sm

    # Build data
    _MG = [
        [("UK", "UK"), ("Brazil", "Brazil"), ("Canada", "Canada ROC"), ("Ontario", "Ontario"), ("Greece", "Greece")],
        [("Southern Europe", "Southern Europe"), ("Austria", "Austria"), ("Denmark", "Denmark"), ("Germany", "Germany"), ("ROW", "ROW")],
        [("Belgium", "Belgium"), ("Italy", "Italy"), ("South Africa", "South Africa"), ("Black Rush", "Black Rush"), ("Foxy NZ", "Foxy NZ")],
    ]

    gg = [[(l, game_counts.get(k, 0)) for l, k in row] for row in _MG]
    if "South Africa" not in opt_counts:
        opt_counts["South Africa"] = 0
    ol = [(m, c) for m, c in opt_counts.items() if c > 0]
    ol.sort(key=lambda x: x[1], reverse=True)
    zero_markets = [(m, c) for m, c in opt_counts.items() if c == 0]
    ol = ol + zero_markets

    tg = total_games_override if total_games_override else sum(c for row in gg for _, c in row if c > 0)
    to = total_opts_override if total_opts_override else sum(c for _, c in ol)
    tc = total_cont_override if total_cont_override else len(cont_list)

    active_cards = []
    for ri, row in enumerate(gg):
        for ci, (lbl, cnt) in enumerate(row):
            if cnt > 0 and lbl:
                clr = CARD_COLORS[ri % 3][ci % 5]
                active_cards.append((lbl, cnt, clr))

    # Calculate total height needed
    rows_needed = (len(active_cards) + 4) // 5
    opt_rows = len(ol) + 1  # +1 for TOTAL
    cont_rows = len(cont_list)
    table_rows = max(opt_rows, cont_rows)
    ku_lines = len(key_updates)

    total_h = (60 +  # top bar + title
               80 +  # KPI boxes
               20 +  # section header
               rows_needed * 70 +  # cards grid
               30 +  # section header
               table_rows * 22 + 40 +  # tables
               30 +  # section header
               ku_lines * 18 + 20 +  # key updates
               40)  # footer
    total_h = max(total_h, 800)

    # Create image
    img = Image.new('RGB', (W, total_h), WHITE_C)
    draw = ImageDraw.Draw(img)
    y = 0

    # Top accent bar
    draw.rectangle([0, 0, W, 6], fill=PURPLE_C)
    y = 20

    # Title
    draw.text((MARGIN, y), "LOBBY-OPS  |  WEEKLY SUMMARY", fill=PURPLE_C, font=font_title)
    y += 30
    draw.text((MARGIN, y), f"{week_label}  |  P&T Global Gaming Content", fill=DARK_C, font=font_sm)
    y += 30

    # KPI Boxes
    kpi_w = (CW - 20) // 3
    kpis = [(tg, "GAMES RELEASED"), (to, "OPTIMIZATIONS"), (tc, "NEW CONTAINERS")]
    for i, (val, lbl) in enumerate(kpis):
        x = MARGIN + i * (kpi_w + 10)
        draw.rounded_rectangle([x, y, x + kpi_w, y + 70], radius=8, fill=KPI_BG_C)
        draw.text((x + kpi_w//2, y + 12), lbl, fill=DARK_C, font=font_tiny, anchor="mt")
        draw.text((x + kpi_w//2, y + 32), str(val), fill=PURPLE_C, font=font_xl, anchor="mt")
    y += 85

    # Games Released section
    draw.text((MARGIN, y), "GAMES RELEASED BY MARKET", fill=PURPLE_C, font=font_section)
    y += 18
    draw.rectangle([MARGIN, y, MARGIN + 200, y + 3], fill=PURPLE_C)
    y += 12

    # Cards grid
    card_w = (CW - 4 * 8) // 5
    card_h = 55
    for idx, (lbl, cnt, clr) in enumerate(active_cards):
        ri = idx // 5
        ci = idx % 5
        cx = MARGIN + ci * (card_w + 8)
        cy = y + ri * (card_h + 8)
        draw.rounded_rectangle([cx, cy, cx + card_w, cy + card_h], radius=6, fill=clr)
        draw.text((cx + card_w//2, cy + 10), str(cnt), fill=DARK_C, font=font_lg, anchor="mt")
        draw.text((cx + card_w//2, cy + 38), lbl, fill=DARK_C, font=font_tiny, anchor="mt")
    y += rows_needed * (card_h + 8) + 15

    # Optimizations table (left - narrow) + Containers table (right - wider)
    left_w = int(CW * 0.33)  # 33% for opt table
    right_w = CW - left_w - 20  # 67% for containers

    # Opt header
    draw.text((MARGIN, y), "OPTIMIZATIONS BY MARKET", fill=PURPLE_C, font=font_section)
    draw.text((MARGIN + left_w + 20, y), "NEW CONTAINERS ENABLED", fill=PURPLE_C, font=font_section)
    y += 18
    draw.rectangle([MARGIN, y, MARGIN + 180, y + 3], fill=PURPLE_C)
    draw.rectangle([MARGIN + left_w + 20, y, MARGIN + left_w + 200, y + 3], fill=PURPLE_C)
    y += 8

    # Table headers
    draw.text((MARGIN + 5, y), "Market", fill=DARK_C, font=font_sm)
    draw.text((MARGIN + left_w - 45, y), "Count", fill=DARK_C, font=font_sm)
    # Container column widths: auto-adjust based on content
    max_cn_len = max((len(cn) for cn, _ in cont_list), default=20) if cont_list else 20
    max_cb_len = max((len(cb) for _, cb in cont_list), default=10) if cont_list else 10
    cn_ratio = max_cn_len / (max_cn_len + max_cb_len)
    cn_col_w = int(right_w * min(cn_ratio, 0.75))  # container name gets up to 75%
    cb_col_w = right_w - cn_col_w
    rx = MARGIN + left_w + 20
    draw.text((rx + 5, y), "Container", fill=DARK_C, font=font_sm)
    draw.text((rx + cn_col_w + 5, y), "Brand", fill=DARK_C, font=font_sm)
    y += 20

    # Data rows
    row_h = 20
    od = ol + [("TOTAL", to)]
    max_rows = max(len(od), len(cont_list))
    for i in range(max_rows):
        ry = y + i * row_h
        # Opt row
        if i < len(od):
            mk, cn = od[i]
            is_total = mk == "TOTAL"
            bg = TOTAL_BG_C if is_total else (ROW1_C if i % 2 == 0 else ROW2_C)
            draw.rectangle([MARGIN, ry, MARGIN + left_w, ry + row_h], fill=bg)
            draw.text((MARGIN + 5, ry + 3), mk, fill=DARK_C, font=font_tiny)
            draw.text((MARGIN + left_w - 45, ry + 3), str(cn), fill=DARK_C, font=font_tiny)

        # Container row
        if i < len(cont_list):
            cn_name, cn_brand = cont_list[i]
            bg = ROW1_C if i % 2 == 0 else ROW2_C
            draw.rectangle([rx, ry, rx + right_w, ry + row_h], fill=bg)
            # Auto-fit: calculate max chars based on column pixel width (~7px per char at font_tiny)
            max_cn_chars = cn_col_w // 7
            max_cb_chars = cb_col_w // 7
            cn_disp = cn_name if len(cn_name) <= max_cn_chars else cn_name[:max_cn_chars - 3] + "..."
            cb_disp = cn_brand if len(cn_brand) <= max_cb_chars else cn_brand[:max_cb_chars - 3] + "..."
            draw.text((rx + 5, ry + 3), cn_disp, fill=DARK_C, font=font_tiny)
            draw.text((rx + cn_col_w + 5, ry + 3), cb_disp, fill=DARK_C, font=font_tiny)

    y += max_rows * row_h + 15

    # Key Updates
    draw.text((MARGIN, y), "KEY UPDATES", fill=PURPLE_C, font=font_section)
    y += 18
    draw.rectangle([MARGIN, y, MARGIN + 180, y + 3], fill=PURPLE_C)
    y += 10
    for u in key_updates:
        txt_line = "• " + u
        if len(txt_line) > 100:
            txt_line = txt_line[:97] + "..."
        draw.text((MARGIN + 5, y), txt_line, fill=DARK_C, font=font_tiny)
        y += 16
    y += 10

    # Footer
    draw.text((W // 2, y), "Entain  |  P&T Global Gaming Content  |  Lobby-Ops", fill=PURPLE_C, font=font_tiny, anchor="mt")
    y += 20
    draw.rectangle([0, y, W, y + 6], fill=PURPLE_C)

    # Crop to actual content height
    img = img.crop((0, 0, W, y + 6))

    # Save to bytes
    buf = _io.BytesIO()
    img.save(buf, format='PNG', dpi=(150, 150))
    buf.seek(0)
    return buf.getvalue()


def _gen_summary_pptx(game_counts, opt_counts, cont_list, key_updates, week_label, total_games_override=None, total_opts_override=None, total_cont_override=None):
    """Generate portrait PPTX summary - exact replica of build_week5_v3.py output."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    _MG = [
        [("UK", "UK"), ("Brazil", "Brazil"), ("Canada", "Canada ROC"), ("Ontario", "Ontario"), ("Greece", "Greece")],
        [("Southern Europe", "Southern Europe"), ("Austria", "Austria"), ("Denmark", "Denmark"), ("Germany", "Germany"), ("ROW", "ROW")],
        [("Belgium", "Belgium"), ("Italy", "Italy"), ("South Africa", "South Africa"), ("Black Rush", "Black Rush"), ("Foxy NZ", "Foxy NZ")],
    ]
    # Build optimizations list from opt_counts - show all markets, sorted descending
    # Ensure South Africa appears even if 0
    if "South Africa" not in opt_counts:
        opt_counts["South Africa"] = 0
    ol = [(m, c) for m, c in opt_counts.items() if c > 0]
    ol.sort(key=lambda x: x[1], reverse=True)  # descending order
    # Also include markets with 0 count at the end
    zero_markets = [(m, c) for m, c in opt_counts.items() if c == 0]
    ol = ol + zero_markets
    CARD_COLORS = [
        [RGBColor(0x6E,0xE7,0xB7),RGBColor(0xFD,0xE6,0x8A),RGBColor(0xFC,0xA5,0xA5),RGBColor(0xFD,0xBA,0x74),RGBColor(0x67,0xE8,0xF9)],
        [RGBColor(0xC4,0xB5,0xFD),RGBColor(0xF9,0xA8,0xD4),RGBColor(0xFC,0xA5,0xA5),RGBColor(0xFD,0xBA,0x74),RGBColor(0x5E,0xEA,0xD4)],
        [RGBColor(0xFD,0xE6,0x8A),RGBColor(0x6E,0xE7,0xB7),RGBColor(0x67,0xE8,0xF9),RGBColor(0xC4,0xB5,0xFD),RGBColor(0xF9,0xA8,0xD4)],
    ]
    PURPLE_A = RGBColor(0x5B,0x21,0xB6)
    BLACK = RGBColor(0x00,0x00,0x00)
    DARK = RGBColor(0x1F,0x1F,0x1F)
    KPI_BG = RGBColor(0xF3,0xF4,0xF6)
    ROW1_C = RGBColor(0xED,0xE9,0xFE)
    ROW2_C = RGBColor(0xF5,0xF3,0xFF)
    TOTAL_BG = RGBColor(0xE8,0xDE,0xF8)

    # Build data
    gg = [[(l, game_counts.get(k, 0)) for l, k in row] for row in _MG]
    # Games Released total = sum of all region cards INCLUDING Black Rush and Foxy NZ
    tg = total_games_override if total_games_override else sum(c for row in gg for lbl, c in row if c > 0 and lbl)
    to = total_opts_override if total_opts_override else sum(c for _, c in ol)
    tc = total_cont_override if total_cont_override else len(cont_list)

    # Create portrait PPTX (7.5" x 13.333")
    prs_out = Presentation()
    prs_out.slide_width = Inches(7.5)
    prs_out.slide_height = Inches(13.333)
    slide = prs_out.slides.add_slide(prs_out.slide_layouts[6])

    SW = prs_out.slide_width
    SH = prs_out.slide_height
    MARGIN = Emu(274320)
    CONTENT_W = SW - 2 * MARGIN

    def rect(l, t, w, h, clr):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

    def txt(l, t, w, h, text, sz=Pt(11), bold=False, clr=DARK, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(27432); tf.margin_right = Emu(27432)
        tf.margin_top = Emu(18288); tf.margin_bottom = Emu(18288)
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text; run.font.name = 'Mulish'
        run.font.size = sz; run.font.bold = bold; run.font.color.rgb = clr

    def card(l, t, w, h, fill, number, label):
        rect(l, t, w, h, fill)
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = str(number); run.font.name = 'Mulish'
        run.font.size = Pt(20); run.font.bold = True; run.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)
        run2 = p2.add_run(); run2.text = label; run2.font.name = 'Mulish'
        run2.font.size = Pt(7); run2.font.color.rgb = DARK

    # Top bar
    rect(0, 0, SW, Emu(54864), PURPLE_A)
    # Title
    txt(MARGIN, Emu(137160), CONTENT_W, Emu(365760), "LOBBY-OPS  |  WEEKLY SUMMARY", sz=Pt(14), bold=True, clr=PURPLE_A)
    txt(MARGIN, Emu(457200), CONTENT_W, Emu(228600), f"{week_label}  |  P&T Global Gaming Content", sz=Pt(9), clr=DARK)

    # KPIs
    kpi_y = Emu(822960); kpi_h = Emu(594360); kpi_w = Emu(2011680); kpi_gap = Emu(137160)
    kpi_x0 = int((SW - 3*kpi_w - 2*kpi_gap) / 2)
    for i, (val, lbl) in enumerate([(tg, "GAMES RELEASED"), (to, "OPTIMIZATIONS"), (tc, "NEW CONTAINERS")]):
        x = kpi_x0 + i*(kpi_w+kpi_gap)
        rect(x, kpi_y, kpi_w, kpi_h, KPI_BG)
        box = slide.shapes.add_textbox(x, kpi_y, kpi_w, kpi_h)
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = lbl; run.font.name = 'Mulish'; run.font.size = Pt(8); run.font.bold = True; run.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run(); run2.text = str(val); run2.font.name = 'Mulish'; run2.font.size = Pt(24); run2.font.bold = True; run2.font.color.rgb = PURPLE_A

    # Games Released section
    sec_y = Emu(1600200)
    txt(MARGIN, sec_y, CONTENT_W, Emu(256032), "GAMES RELEASED BY MARKET", sz=Pt(9), bold=True, clr=PURPLE_A)
    rect(MARGIN, sec_y+Emu(256032), Emu(2286000), Emu(18288), PURPLE_A)

    # Cards grid - only show markets with games > 0, auto-fit layout
    # Flatten to only non-zero markets
    active_cards = []
    for ri, row in enumerate(gg):
        for ci, (lbl, cnt) in enumerate(row):
            if cnt > 0 and lbl:
                clr = CARD_COLORS[ri % len(CARD_COLORS)][ci % 5]
                active_cards.append((lbl, cnt, clr))

    num_cards = len(active_cards)
    cols = 5  # keep 5 columns
    rows_needed = (num_cards + cols - 1) // cols

    cw = Emu(1188720); ch = Emu(457200); cg = Emu(64008)
    cx0 = int((SW - cols*cw - (cols-1)*cg)/2); gy = sec_y + Emu(320000)
    for idx, (lbl, cnt, clr) in enumerate(active_cards):
        ri = idx // cols
        ci = idx % cols
        card(cx0+ci*(cw+cg), gy+ri*(ch+cg), cw, ch, clr, cnt, lbl)

    # Optimizations table - position after cards grid
    opt_y = gy + rows_needed*(ch+cg) + Emu(200000)
    txt(MARGIN, opt_y, Emu(2400000), Emu(256032), "OPTIMIZATIONS BY MARKET", sz=Pt(9), bold=True, clr=PURPLE_A)
    rect(MARGIN, opt_y+Emu(256032), Emu(1800000), Emu(18288), PURPLE_A)
    hdr_y = opt_y + Emu(350000); rh = Emu(155448)
    col1w = Emu(2400000); col2w = Emu(500000)  # reduced opt table width
    rect(MARGIN, hdr_y, col1w, Emu(9144), PURPLE_A)
    txt(MARGIN, hdr_y, Emu(1900000), Emu(182880), "Market", sz=Pt(9), bold=True, clr=BLACK)
    txt(MARGIN+Emu(1900000), hdr_y, col2w, Emu(182880), "Count", sz=Pt(9), bold=True, clr=BLACK, align=PP_ALIGN.CENTER)
    dy = hdr_y + Emu(228600)
    od = ol + [("TOTAL", to)]
    for i, (mk, cn) in enumerate(od):
        y = dy + i*rh; it = (mk == "TOTAL")
        bg = TOTAL_BG if it else (ROW1_C if i%2==0 else ROW2_C)
        rect(MARGIN, y, col1w, rh, bg)
        txt(MARGIN+Emu(45720), y, Emu(1900000), rh, mk, sz=Pt(8), bold=it, clr=BLACK)
        txt(MARGIN+Emu(1900000), y, col2w, rh, str(cn), sz=Pt(8), bold=it, clr=BLACK, align=PP_ALIGN.CENTER)
    rect(MARGIN, dy+len(od)*rh, col1w, Emu(9144), PURPLE_A)

    # Containers table (right side) - wider now
    cont_y = opt_y
    ctx = MARGIN + col1w + Emu(150000)  # smaller gap between tables
    ctw = SW - ctx - MARGIN
    txt(ctx, cont_y, ctw, Emu(256032), "NEW CONTAINERS ENABLED", sz=Pt(9), bold=True, clr=PURPLE_A)
    rect(ctx, cont_y+Emu(256032), Emu(2000000), Emu(18288), PURPLE_A)
    # Header - give more space to container name column
    cnw = int(ctw * 0.70); cbw = ctw - cnw
    rect(ctx, hdr_y, ctw, Emu(9144), PURPLE_A)
    txt(ctx, hdr_y, cnw, Emu(182880), "Container", sz=Pt(9), bold=True, clr=BLACK)
    txt(ctx+cnw, hdr_y, cbw, Emu(182880), "Brand", sz=Pt(9), bold=True, clr=BLACK)
    # Rows - enable word wrap, increase row height to allow text wrapping
    cont_rh = Emu(220000)  # taller rows to allow wrapped text
    for i, (cn, cb) in enumerate(cont_list):
        y = dy + i * cont_rh
        bg = ROW1_C if i%2==0 else ROW2_C
        rect(ctx, y, ctw, cont_rh, bg)
        # No truncation - allow word wrap instead
        cn_display = cn
        cb_display = cb if len(cb) <= 18 else cb[:15] + "..."
        # Create text boxes with word_wrap ENABLED so long names flow to next line
        box_cn = slide.shapes.add_textbox(ctx+Emu(27432), y, cnw-Emu(27432), cont_rh)
        tf_cn = box_cn.text_frame
        tf_cn.word_wrap = True
        tf_cn.auto_size = None
        tf_cn.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_cn.margin_left = Emu(0); tf_cn.margin_top = Emu(0); tf_cn.margin_bottom = Emu(0)
        p_cn = tf_cn.paragraphs[0]; p_cn.alignment = PP_ALIGN.LEFT
        run_cn = p_cn.add_run(); run_cn.text = cn_display
        run_cn.font.name = 'Mulish'; run_cn.font.size = Pt(7); run_cn.font.color.rgb = BLACK

        box_cb = slide.shapes.add_textbox(ctx+cnw, y, cbw, cont_rh)
        tf_cb = box_cb.text_frame
        tf_cb.word_wrap = True
        tf_cb.auto_size = None
        tf_cb.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_cb.margin_left = Emu(0); tf_cb.margin_top = Emu(0); tf_cb.margin_bottom = Emu(0)
        p_cb = tf_cb.paragraphs[0]; p_cb.alignment = PP_ALIGN.LEFT
        run_cb = p_cb.add_run(); run_cb.text = cb_display
        run_cb.font.name = 'Mulish'; run_cb.font.size = Pt(7); run_cb.font.color.rgb = BLACK

    # Key Updates - account for potentially different table heights
    opt_table_end = dy + len(od) * rh
    cont_table_end = dy + len(cont_list) * cont_rh
    ku_y = max(opt_table_end, cont_table_end) + Emu(200000)
    txt(MARGIN, ku_y, CONTENT_W, Emu(256032), "KEY UPDATES", sz=Pt(9), bold=True, clr=PURPLE_A)
    rect(MARGIN, ku_y+Emu(256032), Emu(2286000), Emu(18288), PURPLE_A)
    up_y = ku_y + Emu(320000)
    # Calculate remaining space based on content
    ku_height = len(key_updates) * Emu(140000) + Emu(100000)
    box = slide.shapes.add_textbox(MARGIN, up_y, CONTENT_W, ku_height)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None; tf.margin_left = Emu(36000)
    for i, u in enumerate(key_updates):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(2); p.space_after = Pt(1)
        run = p.add_run(); run.text = f"\u2022 {u}"; run.font.name = 'Mulish'; run.font.size = Pt(6.5); run.font.color.rgb = BLACK

    # Footer - position right after key updates (eliminate white space)
    footer_y = up_y + ku_height + Emu(200000)
    # Resize slide height to fit content tightly
    new_height = footer_y + Emu(350000)
    prs_out.slide_height = new_height
    txt(MARGIN, new_height-Emu(280000), CONTENT_W, Emu(228600), "Entain  |  P&T Global Gaming Content  |  Lobby-Ops", sz=Pt(7), clr=PURPLE_A, align=PP_ALIGN.CENTER)
    rect(0, new_height-Emu(54864), SW, Emu(54864), PURPLE_A)

    # Save to bytes
    buf = io.BytesIO(); prs_out.save(buf); buf.seek(0)
    return buf.getvalue()


# ============ TAB 2: Quick Summary (Excel source + editable key updates) ============
with tab2:
    st.subheader("\U0001f4dd Quick Summary")
    st.markdown("Upload **Major Updates.xlsx** and optionally **W2 CSV** \u2192 auto-generate summary report (PPTX for PDF export)")

    col_t3a, col_t3b = st.columns(2)
    with col_t3a:
        t3_file = st.file_uploader("Upload Major Updates.xlsx", type=["xlsx"], key="t3_xlsx")
    with col_t3b:
        t3_csv = st.file_uploader("Upload W2 CSV (optional, for game counts)", type=["csv"], key="t3_csv")
    t3_week = st.text_input("Week Label", value="Jun Week 2", key="t3_wk")

    # Step 1: Parse data when file is uploaded
    if t3_file:
        file_key = f"{t3_file.name}_{t3_file.size}_v5"
        if st.session_state.get('t3_fkey') != file_key:
            # Clear old cached data
            for k in ['t3_gc', 't3_oc', 't3_cr', 't3_ku', 't3_fkey']:
                st.session_state.pop(k, None)
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx")
            tmp.write(t3_file.read()); tmp.close()
            t3_file.seek(0)
            try:
                gc, oc, cr, ku = parse_excel_for_summary(tmp.name)

                # If CSV uploaded, use for game counts
                if t3_csv:
                    csv_content = t3_csv.read().decode("utf-8")
                    t3_csv.seek(0)
                    csv_games = parse_csv(csv_content)
                    csv_gc = {}
                    for g in csv_games:
                        sk = MARKET_TO_SLIDE.get(g["market"])
                        if sk:
                            summary_map = {"UK": "UK", "Brazil": "Brazil", "Ontario": "Ontario",
                                "ROC": "Canada ROC", "ROW": "ROW", "Greece": "Greece",
                                "Southern Europe": "Southern Europe", "Austria": "Austria",
                                "Denmark": "Denmark", "Germany": "Germany",
                                "Belgium": "Belgium", "Italy": "Italy", "South Africa": "South Africa"}
                            sum_key = summary_map.get(sk, sk)
                            csv_gc[sum_key] = csv_gc.get(sum_key, 0) + 1
                    for k, v in csv_gc.items():
                        if v > 0: gc[k] = v

                # Clean containers
                raw_containers = {}
                for desc, brand in cr:
                    nm = desc.strip()
                    for pfx in ["\u2022 ", "\u2022", "• "]:
                        if nm.startswith(pfx): nm = nm[len(pfx):]
                    nm = nm.strip()
                    if not nm or nm == "NA" or len(nm) < 3:
                        continue
                    # Try to extract quoted container name first
                    quoted = re.findall(r'["\u201c\u201d\u2018\u2019\'"]([^"\u201c\u201d\u2018\u2019\'"]+)["\u201c\u201d\u2018\u2019\'"]', nm)
                    if quoted:
                        container_name = quoted[0].strip()
                    else:
                        # Extract short name from description
                        container_name = nm
                        # Remove common prefixes to get to the name
                        for pfx in ["New Provider ", "New provider ", "New Container ", "New container ",
                                     "Added hidden container ", "Added new container ",
                                     "Added casino hidden container ", "Enabled container ",
                                     "Enabled new container ", "Created container ",
                                     "New hidden container "]:
                            if container_name.startswith(pfx):
                                container_name = container_name[len(pfx):]
                                break
                        # Cut at common stop phrases to isolate the container name
                        for stop in [" container has been", " has been", " new container",
                                     " was ", " enabled for ", " enabled on ", " enabled in ",
                                     " added to ", " added for ", " added on ",
                                     " for Bwin", " for Premium", " for Party", " for Ladbrokes",
                                     " for Coral", " for Gala", " for Foxy", " for Sportingbet",
                                     " for all ", " across ",
                                     " - Bwin", " - Premium", " - Party", " - Ladbrokes",
                                     " - Coral", " - Gala", " - Foxy", " - Sportingbet",
                                     " in ", " on "]:
                            idx = container_name.lower().find(stop.lower())
                            if idx > 3:
                                container_name = container_name[:idx]
                                break
                        container_name = container_name.strip().strip("'\".,;:").strip()
                    
                    # Final cleanup: remove any trailing brand suffixes like "- Bwin.Com, Premium.Com"
                    # Pattern: name followed by " - " and domain-like text
                    brand_suffix = re.search(r'\s*[-–]\s*(\w+\.\w+.*)', container_name)
                    if brand_suffix and brand_suffix.start() > 3:
                        container_name = container_name[:brand_suffix.start()].strip()

                    if container_name and len(container_name) > 2:
                        # Split on & or comma to count each container separately
                        # e.g. "Top Bwin Originals & More Bwin Originals" → 2 containers
                        parts = re.split(r'\s*&\s*|\s*,\s*', container_name)
                        # Filter out parts that look like brand names (contain dots like Bwin.com)
                        container_parts = [p.strip() for p in parts if p.strip() and len(p.strip()) > 2 and '.' not in p.strip()]
                        if not container_parts:
                            container_parts = [container_name]
                        
                        # Extract brand name from the sentence - match known brand patterns
                        brand_patterns = [
                            r'(Ladbrokes\.\w+)', r'(Coral\.\w+)',
                            r'(Bwin\.\w+)', r'(Bwincasino\.\w+)', r'(Bwindice\.\w+)',
                            r'(Sportingbet\.\w+)', r'(Betboo\.\w+)',
                            r'(Partycasino\.\w+)', r'(Partypoker\.\w+)',
                            r'(Premium\.\w+)', r'(Giocodigitale\.\w+)',
                            r'(Vistabet\.\w+)', r'(GD\.\w+)',
                            r'(Gala Casino)', r'(Gala)', r'(Foxy)',
                        ]
                        # Also try to match "for/on [BrandLabel]" at end of sentence
                        brand_label_patterns = [
                            r'\bfor\s+((?:Coral|Ladbrokes|Bwin|Sportingbet|Partycasino|Partypoker|Premium|Gala|Foxy|Betboo|Vistabet|GD|Giocodigitale)[\w.]*(?:\s*(?:&|,)\s*(?:Coral|Ladbrokes|Bwin|Sportingbet|Partycasino|Partypoker|Premium|Gala|Foxy|Betboo|Vistabet|GD|Giocodigitale)[\w.]*)*)',
                            r'\bon\s+((?:Coral|Ladbrokes|Bwin|Sportingbet|Partycasino|Partypoker|Premium|Gala|Foxy|Betboo|Vistabet|GD|Giocodigitale)[\w.]*(?:\s*(?:&|,)\s*(?:Coral|Ladbrokes|Bwin|Sportingbet|Partycasino|Partypoker|Premium|Gala|Foxy|Betboo|Vistabet|GD|Giocodigitale)[\w.]*)*)',
                            r'\bon\s+(Spain|Italy|Greece|Portugal|Belgium|Germany|Denmark|Austria|Brazil|South Africa)',
                            r'\bfor\s+(Spain|Italy|Greece|Portugal|Belgium|Germany|Denmark|Austria|Brazil|South Africa)',
                        ]
                        brand_found = ""
                        # First try label patterns (more specific - "for Coral.uk & Ladbrokes.com")
                        for bp in brand_label_patterns:
                            m = re.search(bp, desc, re.IGNORECASE)
                            if m:
                                brand_found = m.group(1).strip()
                                break
                        # Then try individual brand patterns
                        if not brand_found:
                            for bp in brand_patterns:
                                m = re.search(bp, desc, re.IGNORECASE)
                                if m:
                                    brand_found = m.group(1).strip()
                                    break
                        if not brand_found:
                            brand_found = brand  # fallback to Sheet 2 col A
                        
                        for cp in container_parts:
                            if cp not in raw_containers:
                                raw_containers[cp] = set()
                            raw_containers[cp].add(brand_found)
                clean_cr = [(name.title(), ", ".join(sorted(brands))) for name, brands in raw_containers.items()]

                # Consolidate key updates
                seen = set(); raw_ku = []
                for u in ku:
                    u_clean = u.strip().replace("has been ", "was ").replace("have been ", "were ")
                    if u_clean and u_clean[:40].lower() not in seen:
                        seen.add(u_clean[:40].lower()); raw_ku.append(u_clean)
                consolidated = []
                eznav_b = []; page_b = []; thumb_b = []; other_b = []
                for u in raw_ku:
                    ul = u.lower()
                    if "eznav" in ul and "optimization" in ul: eznav_b.append(u)
                    elif ("home" in ul or "slots" in ul or "new page" in ul) and "optimization" in ul: page_b.append(u)
                    elif "thumbnail" in ul: thumb_b.append(u)
                    else: other_b.append(u)
                if page_b: consolidated.append("Page optimizations were completed across all regions (Home, Slots, New Games, Trending, Exclusives, Live Casino).")
                if eznav_b: consolidated.append("Eznav optimization was completed for all labels across all regions.")
                if thumb_b: consolidated.append("Game Thumbnails were updated across applicable labels.")
                for u in other_b:
                    if len(consolidated) >= 10: break
                    consolidated.append(u.strip().rstrip('.') + '.')

                st.session_state['t3_gc'] = gc
                st.session_state['t3_oc'] = oc
                st.session_state['t3_cr'] = clean_cr
                st.session_state['t3_ku'] = consolidated
                st.session_state['t3_fkey'] = file_key
                os.unlink(tmp.name)
            except Exception as e:
                st.error(f"Error parsing: {e}")
                import traceback; st.code(traceback.format_exc())

    # Step 2: Show editable key updates + generate
    if 't3_gc' in st.session_state and st.session_state.get('t3_fkey'):
        gc = st.session_state['t3_gc']
        oc = st.session_state['t3_oc']
        clean_cr = st.session_state['t3_cr']

        total_g = sum(gc.values())
        total_o = sum(oc.values())
        st.success(f"\u2705 Games: {total_g} | Optimizations: {total_o} | Containers: {len(clean_cr)}")

        st.subheader("Key Updates")
        st.caption("Auto-picked from file. Edit, add or remove lines below.")
        ku_default = "\n".join(st.session_state['t3_ku'])
        ed_ku = st.text_area("Key Updates (one per line)", value=ku_default, height=200, key="t3_ku_edit")

        if st.button("\U0001f680 Generate Summary", type="primary", use_container_width=True, key="t3_gen"):
            final_ku = [l.strip() for l in ed_ku.split("\n") if l.strip()]

            with st.spinner("Generating..."):
                try:
                    pptx_bytes = _gen_summary_pptx(gc, oc, clean_cr, final_ku, t3_week, total_g, total_o, len(clean_cr))
                    png_bytes = _gen_summary_png(gc, oc, clean_cr, final_ku, t3_week, total_g, total_o, len(clean_cr))

                    st.divider()

                    c1, c2 = st.columns(2)
                    with c1:
                        st.download_button("\U0001f4e5 Download PPTX", pptx_bytes,
                            file_name=f"LobbyOps_Summary_{t3_week.replace(' ','_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True)
                    with c2:
                        st.download_button("\U0001f5bc\ufe0f Download PNG", png_bytes,
                            file_name=f"LobbyOps_Summary_{t3_week.replace(' ','_')}.png",
                            mime="image/png", use_container_width=True)

                    st.image(png_bytes, caption="Preview", use_container_width=True)

                except Exception as e:
                    st.error(f"Error: {e}")
                    import traceback; st.code(traceback.format_exc())
    elif not t3_file:
        st.info("\U0001f4c1 Upload Major Updates Excel to get started.")
